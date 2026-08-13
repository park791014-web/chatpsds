import streamlit as st
import os
import json
import re
import html
import io
import zipfile
import xml.etree.ElementTree as ET
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
from openai import OpenAI
from utils.ui import hide_streamlit_chrome, paste_listener

# v15: 전체 이미지/PDF 페이지 분석 + 구조화 표 학생 식별 오인식 방지
# 1. 페이지 기본 설정 (사이드바 기본 열림 상태로 고정)
st.set_page_config(page_title="Chat PSDongSung", layout="wide", initial_sidebar_state="expanded")

# Streamlit UI 크롬 정밀 제거
hide_streamlit_chrome()

# ==========================================
# Streamlit 네이티브 높이 스크롤 컨테이너 CSS
# ==========================================
st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;800&family=Noto+Sans+KR:wght@400;500;700&display=swap');
    
    html, body, [class*="css"] {
        font-family: 'Inter', 'Noto Sans KR', sans-serif;
    }

    /* 1. 최외곽 스크롤 차단 및 Flex 구조 정의 */
    [data-testid="stMain"],
    .stMain,
    [data-testid="stAppViewContainer"] .main,
    [data-testid="stAppViewContainer"] .stMain {
        height: 100dvh !important;
        max-height: 100dvh !important;
        overflow-y: hidden !important;
        overflow-x: hidden !important;
        display: flex !important;
        flex-direction: column !important;
        min-height: 0 !important;
        overscroll-behavior: none !important;
    }

    /* 2. 메인 콘텐츠 영역 (stMain 내부에서 남은 공간을 Flex로 공유) */
    [data-testid="stMain"] .block-container,
    .block-container {
        flex: 1 1 auto !important;
        min-height: 0 !important;
        height: auto !important;
        max-height: none !important;
        padding-top: 3.2rem !important;
        padding-bottom: 0.5rem !important;
        box-sizing: border-box !important;
        display: flex !important;
        flex-direction: column !important;
    }

    /* 1080px급 이하 화면 (100% 표준 배율 대응) 상단 여백 최소화 */
    @media (max-height: 1100px) {
        [data-testid="stMain"] .block-container,
        .block-container {
            padding-top: 2.5rem !important;
        }
    }

    /* 3. 상단 top_area (수축 방지 및 일반 플로우 아이템 배치) */
    .st-key-top_area,
    div[key="top_area"] {
        flex: 0 0 auto !important;
        position: relative !important;
        top: auto !important;
        z-index: 10 !important;
        background-color: #0f172a !important;
        padding-top: 0.1rem !important;
        padding-bottom: 0.2rem !important;
        border-bottom: none !important;
        margin-bottom: 0.4rem !important;
        overflow: visible !important;
    }

    /* 3.1 top_area 내부 부모 clipping 해제 (100% 배율 로고 잘림 원천 해결) */
    .st-key-top_area [data-testid="stMarkdownContainer"],
    div[key="top_area"] [data-testid="stMarkdownContainer"] {
        overflow: visible !important;
    }

    .st-key-top_area [data-testid="stColumn"],
    div[key="top_area"] [data-testid="stColumn"],
    .st-key-top_area .stColumn,
    div[key="top_area"] .stColumn {
        overflow: visible !important;
    }

    /* 4. 화면 높이에 따라 자동으로 늘고 줄어드는 본문 스크롤 영역 */
    .st-key-content_area,
    div[key="content_area"] {
        height: calc(100dvh - 13.5rem) !important;
        max-height: calc(100dvh - 13.5rem) !important;
        min-height: 260px !important;
        overflow-y: auto !important;
        overflow-x: hidden !important;
        box-sizing: border-box !important;
        padding-bottom: 3rem !important;
    }

    /* 일반 챗봇은 고정 690px 대신 실제 뷰포트의 남은 높이를 사용 */
    .st-key-chat_history_area,
    div[key="chat_history_area"] {
        height: calc(100dvh - 15.5rem) !important;
        max-height: calc(100dvh - 15.5rem) !important;
        min-height: 260px !important;
        overflow-y: auto !important;
        overflow-x: hidden !important;
        box-sizing: border-box !important;
        padding-bottom: 1.5rem !important;
        background: transparent !important;
    }

    @media (max-height: 760px) {
        .st-key-chat_history_area, div[key="chat_history_area"] {
            height: calc(100dvh - 13.8rem) !important;
            max-height: calc(100dvh - 13.8rem) !important;
            min-height: 220px !important;
        }
        .st-key-content_area, div[key="content_area"] {
            height: calc(100dvh - 12.5rem) !important;
            max-height: calc(100dvh - 12.5rem) !important;
            min-height: 220px !important;
        }
    }

    /* 4.1 일반 챗봇 전용 메시지 카드 초슬림화 */
    .st-key-chat_history_area [data-testid="stChatMessage"],
    div[key="chat_history_area"] [data-testid="stChatMessage"] {
        padding-top: 0.28rem !important;
        padding-bottom: 0.28rem !important;
        margin-bottom: 0.2rem !important;
        border-radius: 8px !important;
    }

    .st-key-chat_history_area [data-testid="stChatMessage"] p,
    div[key="chat_history_area"] [data-testid="stChatMessage"] p {
        margin-top: 0 !important;
        margin-bottom: 0.1rem !important;
        line-height: 1.38 !important;
        font-size: 0.94rem !important;
    }

    .st-key-chat_history_area [data-testid="stCaptionContainer"],
    div[key="chat_history_area"] [data-testid="stCaptionContainer"] {
        margin-top: 0.02rem !important;
        margin-bottom: 0 !important;
    }


    /* 마지막 메시지/토큰 표시가 하단 입력창에 가려지지 않도록 안전 여백 확보 */
    .chat-bottom-safe-space {
        height: 96px;
        min-height: 96px;
        width: 100%;
        pointer-events: none;
    }

    /* 5. Streamlit 공식 하단 고정 컨테이너 (st.bottom) 수축 방지 */
    [data-testid="stBottom"],
    div[data-testid="stBottom"],
    [data-testid="stBottomBlockContainer"],
    div[data-testid="stBottomBlockContainer"] {
        flex: 0 0 auto !important;
        padding-top: 0 !important;
        padding-bottom: 10px !important;
        margin-top: 0 !important;
        margin-bottom: 0 !important;
        background: transparent !important;
        overflow: visible !important;
        z-index: 1000 !important;
    }

    /* Streamlit이 하단 고정 입력창 높이를 레이아웃에 반영하도록 기본 spacer를 유지 */
    [data-testid="stBottomSpace"],
    div[data-testid="stBottomSpace"] {
        min-height: 64px !important;
    }

    /* 6. 통합 ChatGPT 스타일 Composer 외곽 카드 비활성화 (투명화) */
    .st-key-chat_composer_area,
    div[key="chat_composer_area"] {
        background-color: transparent !important;
        border: none !important;
        border-radius: 0 !important;
        padding: 0 !important;
        margin-bottom: 0 !important;
        box-sizing: border-box !important;
        overflow: visible !important;
        box-shadow: none !important;
    }

    /* 대기 첨부 칩스 패널 (테두리 및 배경 제거) */
    div[key="chat_pending_chips"] {
        background-color: transparent !important;
        padding: 4px 0 !important;
        margin-bottom: 6px !important;
        border-bottom: none !important;
    }

    /* ChatGPT 스타일 콤팩트 썸네일 (최대 80x60px) */
    div[key="chat_pending_chips"] img,
    .st-key-chat_pending_chips img {
        max-width: 80px !important;
        max-height: 60px !important;
        width: auto !important;
        height: auto !important;
        object-fit: contain !important;
        border-radius: 6px !important;
        border: 1px solid #475569 !important;
        display: block !important;
        margin-bottom: 2px !important;
    }

    /* 칩스 내 삭제 버튼 (×) 콤팩트화 */
    div[key="chat_pending_chips"] button {
        height: 22px !important;
        min-height: 22px !important;
        padding: 0 6px !important;
        font-size: 0.8rem !important;
        line-height: 1 !important;
        border-radius: 4px !important;
        background-color: #334155 !important;
        color: #94a3b8 !important;
        border: none !important;
    }
    div[key="chat_pending_chips"] button:hover {
        background-color: #ef4444 !important;
        color: #ffffff !important;
    }

    /* st.chat_input 자체를 1개짜리 단일 테두리 박스 디자인으로 적용 */
    div[key="chat_composer_area"] [data-testid="stChatInput"],
    div[key="chat_composer_area"] div.stChatInput {
        background-color: #1e293b !important;
        border: 1px solid #38bdf8 !important;
        border-radius: 14px !important;
        margin: 0 !important;
        padding: 4px 8px !important;
        min-height: 38px !important;
        height: 38px !important;
        box-sizing: border-box !important;
    }

    div[key="chat_composer_area"] [data-testid="stChatInput"] [data-baseweb="textarea"],
    div[key="chat_composer_area"] div.stChatInput [data-baseweb="textarea"],
    div[key="chat_composer_area"] [data-testid="stChatInput"] [data-baseweb="base-input"] {
        min-height: 32px !important;
        max-height: 32px !important;
        height: 32px !important;
        margin: 0 !important;
        padding: 0 !important;
        background: transparent !important;
        border: none !important;
        box-shadow: none !important;
    }

    div[key="chat_composer_area"] [data-testid="stChatInput"] textarea,
    div[key="chat_composer_area"] div.stChatInput textarea {
        min-height: 32px !important;
        max-height: 32px !important;
        height: 32px !important;
        line-height: 22px !important;
        font-size: 0.88rem !important;
        padding-top: 5px !important;
        padding-bottom: 5px !important;
        padding-left: 10px !important;
        resize: none !important;
        overflow-y: hidden !important;
        background: transparent !important;
        border: none !important;
        box-sizing: border-box !important;
    }

    div[key="chat_composer_area"] [data-testid="stChatInput"] button {
        align-self: center !important;
        margin-top: auto !important;
        margin-bottom: auto !important;
        height: 26px !important;
        min-height: 26px !important;
        max-height: 26px !important;
        width: 26px !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
    }

    /* 일반 챗봇 대기 첨부자료 칩스 패널 */
    div[key="chat_pending_chips"] {
        background-color: #0f172a !important;
        padding: 0.5rem 0.8rem !important;
        border-radius: 10px !important;
        border: 1px solid #334155 !important;
        margin-top: 0.6rem !important;
        margin-bottom: 0.6rem !important;
    }
    div[key="chat_pending_chips"] img {
        max-width: 130px !important;
        max-height: 95px !important;
        width: auto !important;
        height: auto !important;
        object-fit: contain !important;
        border-radius: 6px !important;
        border: 1px solid #475569 !important;
        display: block !important;
        margin-bottom: 0.3rem !important;
    }

    /* 채팅 히스토리 메시지 내 이미지 크기 제한 및 둥근 모서리 */
    .st-key-chat_history_area [data-testid="stChatMessage"] img,
    div[key="chat_history_area"] [data-testid="stChatMessage"] img {
        max-width: 360px !important;
        max-height: 300px !important;
        width: auto !important;
        height: auto !important;
        object-fit: contain !important;
        border-radius: 8px !important;
        border: 1px solid #334155 !important;
    }
    
    /* 메인 화면 브랜드 헤더 여백 및 100% 배율 잘림 보정 */
    .brand-header {
        margin-top: 0 !important;
        margin-bottom: 0.1rem;
        padding-top: 0.1rem;
        padding-bottom: 0.1rem;
    }
    
    .st-key-top_area [data-testid="stSelectbox"],
    div[key="top_area"] [data-testid="stSelectbox"] {
        margin-top: 0 !important;
    }
    
    /* 파일 업로드 드롭존 컴팩트화 (세로 여유 확보) */
    [data-testid="stFileUploader"] section {
        padding: 0.4rem 0.8rem !important;
        min-height: 50px !important;
    }
    [data-testid="stFileUploaderDropzoneInstructions"] {
        margin-top: 0 !important;
        font-size: 0.8rem !important;
    }


    /* AI가 Markdown 제목(#, ## 등)을 반환해도 세 모드에서 글자 크기를 일정하게 유지 */
    .st-key-chat_history_area [data-testid="stMarkdownContainer"] h1,
    .st-key-chat_history_area [data-testid="stMarkdownContainer"] h2,
    .st-key-chat_history_area [data-testid="stMarkdownContainer"] h3,
    .st-key-chat_history_area [data-testid="stMarkdownContainer"] h4,
    .st-key-chat_history_area [data-testid="stMarkdownContainer"] h5,
    .st-key-chat_history_area [data-testid="stMarkdownContainer"] h6,
    .st-key-chat_history_area [data-testid="stMarkdownContainer"] p,
    .st-key-chat_history_area [data-testid="stMarkdownContainer"] li,
    .st-key-draft_result [data-testid="stMarkdownContainer"] h1,
    .st-key-draft_result [data-testid="stMarkdownContainer"] h2,
    .st-key-draft_result [data-testid="stMarkdownContainer"] h3,
    .st-key-draft_result [data-testid="stMarkdownContainer"] h4,
    .st-key-draft_result [data-testid="stMarkdownContainer"] h5,
    .st-key-draft_result [data-testid="stMarkdownContainer"] h6,
    .st-key-draft_result [data-testid="stMarkdownContainer"] p,
    .st-key-draft_result [data-testid="stMarkdownContainer"] li,
    .st-key-eval_result [data-testid="stMarkdownContainer"] h1,
    .st-key-eval_result [data-testid="stMarkdownContainer"] h2,
    .st-key-eval_result [data-testid="stMarkdownContainer"] h3,
    .st-key-eval_result [data-testid="stMarkdownContainer"] h4,
    .st-key-eval_result [data-testid="stMarkdownContainer"] h5,
    .st-key-eval_result [data-testid="stMarkdownContainer"] h6,
    .st-key-eval_result [data-testid="stMarkdownContainer"] p,
    .st-key-eval_result [data-testid="stMarkdownContainer"] li {
        font-size: 0.94rem !important;
        line-height: 1.5 !important;
    }

    .st-key-chat_history_area [data-testid="stMarkdownContainer"] h1,
    .st-key-chat_history_area [data-testid="stMarkdownContainer"] h2,
    .st-key-chat_history_area [data-testid="stMarkdownContainer"] h3,
    .st-key-draft_result [data-testid="stMarkdownContainer"] h1,
    .st-key-draft_result [data-testid="stMarkdownContainer"] h2,
    .st-key-draft_result [data-testid="stMarkdownContainer"] h3,
    .st-key-eval_result [data-testid="stMarkdownContainer"] h1,
    .st-key-eval_result [data-testid="stMarkdownContainer"] h2,
    .st-key-eval_result [data-testid="stMarkdownContainer"] h3 {
        font-weight: 700 !important;
        margin-top: 0.65rem !important;
        margin-bottom: 0.25rem !important;
    }

    /* draft_result 내부 초안 편집 textarea: 비교/검수 공간 확보를 위해 콤팩트 높이 */
    .st-key-draft_result textarea,
    div[key="draft_result"] textarea {
        min-height: 112px !important;
        height: 112px !important;
        display: block !important;
        visibility: visible !important;
        opacity: 1 !important;
    }

    /* 생기부 작성/검수 전용 컴팩트 타이포그래피 */
    .st-key-eval_result h1, .st-key-eval_result h2, .st-key-eval_result h3,
    div[key="eval_result"] h1, div[key="eval_result"] h2, div[key="eval_result"] h3,
    .st-key-draft_result h1, .st-key-draft_result h2, .st-key-draft_result h3,
    div[key="draft_result"] h1, div[key="draft_result"] h2, div[key="draft_result"] h3 {
        font-size: 0.94rem !important;
        line-height: 1.5 !important;
        margin-top: 0.8rem !important;
        margin-bottom: 0.3rem !important;
        font-weight: 700 !important;
        color: #38bdf8 !important;
    }
    
    /* 테마 고정에 따라 항상 밝고 선명한 텍스트 가시성 유지 */
    .brand-title {
        color: #f8fafc !important;
        font-size: 1.6rem;
        font-weight: 800;
        letter-spacing: -0.5px;
        margin: 0;
        line-height: 1.4;
        display: inline-block;
        transform: translateY(4px) !important;
    }
    .brand-accent {
        color: #38bdf8 !important;
        font-weight: 800;
    }
    .brand-sub {
        color: #94a3b8 !important;
        font-size: 0.85rem;
        font-weight: 400;
        margin-top: 0.2rem;
        letter-spacing: -0.2px;
    }
    .sidebar-brand-title {
        color: #f8fafc !important;
        font-size: 1.4rem;
        font-weight: 800;
        letter-spacing: -0.5px;
        text-align: center;
        margin-top: -0.5rem;
        margin-bottom: 0.8rem;
        padding-bottom: 0.6rem;
        border-bottom: 1px solid #334155;
    }

    /* 영화/드라마 보안 시스템 스타일 로그인 카드 (항상 다크 테마 고정) */
    .login-card {
        text-align: center;
    }
    .login-title {
        color: #f8fafc !important;
        font-size: 1.6rem;
        font-weight: 800;
        letter-spacing: -0.5px;
        margin-bottom: 0.3rem;
        line-height: 1.3;
    }
    .login-sub {
        color: #64748b !important;
        font-size: 0.85rem;
        margin-bottom: 1.5rem;
    }
    </style>
""", unsafe_allow_html=True)

# ==========================================
# Phase 5 유틸리티: NEIS 바이트 & 토큰 계산
# ==========================================
def calculate_neis_bytes(text):
    if not text:
        return 0, 0
    
    char_count = len(text)
    byte_count = 0
    
    for char in text:
        if char == '\n':
            byte_count += 2
        elif ord(char) > 127:
            byte_count += 3
        else:
            byte_count += 1
            
    return char_count, byte_count

def get_scope_document_text(scope="student"):
    parts = []
    attachments = st.session_state.get("attachments", {}).get(scope, [])
    for item in attachments:
        if item.get("type") == "document" and item.get("data"):
            parts.append(str(item.get("data", "")))
    return "\n".join(parts)

def infer_student_identity(text):
    text = text or ""
    sid = ""
    name = ""
    m = re.search(r"(?:학번\s*[:：]?\s*)?(\d{5,8})\s*[,/|\- ]+\s*(?:이름\s*[:：]?\s*)?([가-힣]{2,5})", text)
    if m:
        sid, name = m.group(1), m.group(2)
    else:
        m = re.search(r"(?:이름\s*[:：]?\s*)([가-힣]{2,5}).{0,30}?(?:학번\s*[:：]?\s*)(\d{5,8})", text, re.S)
        if m:
            name, sid = m.group(1), m.group(2)
    if not sid:
        m = re.search(r"학번\s*[:：]?\s*(\d{4,8})", text)
        if m:
            sid = m.group(1)
    if not name:
        m = re.search(r"이름\s*[:：]?\s*([가-힣]{2,5})", text)
        if m:
            name = m.group(1)
    return sid, name

def clean_student_name(name, sid=""):
    """학번이 이름으로 중복 표시되는 오인식을 방지합니다."""
    value = str(name or "").strip()
    sid = str(sid or "").strip()
    if not value or value == sid or value.isdigit():
        return ""
    if not re.fullmatch(r"[가-힣]{2,5}", value):
        return ""
    return value

def extract_student_candidates(text):
    text = text or ""
    found = []
    patterns = [
        r"(?:학번\s*[:：]?\s*)?(\d{5,8})\s*[,/|\- ]+\s*(?:이름\s*[:：]?\s*)?([가-힣]{2,5})",
        r"학번\s*[:：]?\s*(\d{4,8}).{0,50}?이름\s*[:：]?\s*([가-힣]{2,5})",
    ]
    for pattern in patterns:
        for m in re.finditer(pattern, text, re.S):
            pair = (m.group(1), m.group(2))
            if pair not in found:
                found.append(pair)
    return found[:50]

def build_student_source_signature(global_hashes=None):
    """현재 학생 첨부 + 공통참조의 원본 기준 서명. PDF 페이지 파생 이미지는 제외합니다."""
    keys = []
    for item in st.session_state.get("attachments", {}).get("student", []):
        # PDF에서 파생된 숨김 페이지 이미지는 원본 문서 hash로 대표하므로 제외
        if item.get("hidden") and item.get("source") in ("pdf_page", "global_pdf_page"):
            continue
        h = item.get("hash") or ""
        if h:
            keys.append(f"student:{h}")
    for h in (global_hashes or []):
        if h:
            keys.append(f"global:{h}")
    return hashlib.sha256("|".join(sorted(set(keys))).encode("utf-8")).hexdigest() if keys else ""


def build_structured_student_analysis_from_text(text):
    """학번/이름 열이 있는 XLSX/CSV는 AI 호출 없이 학생별 행을 직접 분리합니다."""
    marker = "[STRUCTURED_TABLE_TSV]\n"
    if not text or marker not in text:
        return None
    students = []
    for chunk in text.split(marker)[1:]:
        # 다음 첨부문서 블록이 시작되기 전까지만 현재 표로 간주
        chunk = re.split(r"\n\n\[첨부문서:", chunk, maxsplit=1)[0].strip()
        if not chunk:
            continue
        try:
            df = pd.read_csv(io.StringIO(chunk), sep="\t", dtype=str).fillna("")
        except Exception:
            continue
        if "학번" not in df.columns or "이름" not in df.columns:
            continue
        for _, row in df.iterrows():
            sid = str(row.get("학번", "") or "").strip()
            name = clean_student_name(row.get("이름", ""), sid)
            if not sid and not name:
                continue
            details = []
            for col in df.columns:
                if col in ("학번", "이름"):
                    continue
                value = str(row.get(col, "") or "").strip()
                if value and value.lower() != "nan":
                    details.append(f"{col}: {value}")
            summary = " | ".join(details)[:4500]
            students.append({"id": sid, "name": name, "summary": summary})
    if not students:
        return None
    return {"students": students[:60], "common_summary": "", "source_summary": ""}


def normalize_source_analysis(raw_text):
    """모델의 JSON 응답을 안전하게 파싱하고 학생별 요약 구조로 정규화합니다."""
    raw = (raw_text or "").strip()
    raw = re.sub(r"^```(?:json)?\s*|\s*```$", "", raw, flags=re.I | re.S).strip()
    data = json.loads(raw)
    if not isinstance(data, dict):
        raise ValueError("자료 분석 결과가 JSON 객체 형식이 아닙니다.")
    students = []
    for item in data.get("students", []) or []:
        if not isinstance(item, dict):
            continue
        sid = str(item.get("id", "") or "").strip()
        name = clean_student_name(item.get("name", ""), sid)
        summary = str(item.get("summary", "") or "").strip()
        if sid or name or summary:
            students.append({"id": sid, "name": name, "summary": summary[:5000]})
    return {
        "students": students[:60],
        "common_summary": str(data.get("common_summary", "") or "").strip()[:6000],
        "source_summary": str(data.get("source_summary", "") or "").strip()[:6000],
    }


def analyze_student_sources_once(client, model, selected_model_name, global_ref_text="", candidate_hint=""):
    """원본 첨부를 단 한 번 Vision/텍스트로 분석해 이후 재사용할 가벼운 요약을 만듭니다."""
    prompt = (
        "첨부된 학생 자료와 공통 참고자료를 학교생활기록부 작성용으로 한 번만 분석하세요. "
        "이미지형 PDF와 첨부 이미지는 제공된 모든 페이지/이미지를 빠짐없이 확인하세요. "
        "슬라이드의 글자, 표, 도식, 캡션, 작은 문구까지 보이는 범위에서 최대한 읽고, 오타나 비표준 표현도 임의로 고치지 말고 자료에 보이는 의미를 우선 파악하세요. "
        "학생별 자료가 여러 명이면 절대 서로 섞지 말고 학번과 이름 기준으로 분리하세요. "
        "표지, 감사 인사, 출처 목록처럼 학생 활동 근거가 아닌 내용은 요약에서 제외하세요. "
        "추측하거나 새로운 사실을 만들지 마세요. 학생 한 명의 summary는 핵심 근거만 250~400자 정도로 압축하세요.\n"
        f"식별 참고: {candidate_hint or '첨부자료에서 직접 식별'}\n"
        "반드시 아래 JSON 객체 하나만 출력하세요.\n"
        '{"students":[{"id":"학번","name":"이름","summary":"탐구 주제, 과정, 사용 도구나 개념, 결과, 관찰 가능한 역량을 간결히 요약"}],'
        '"common_summary":"모든 학생에게 공통 적용할 성취기준, 평가기준, 작성규칙의 핵심만 요약",'
        '"source_summary":"학생 구분이 어려운 자료의 공통 핵심내용"}'
    )
    if global_ref_text:
        # 텍스트형 공통자료는 최초 분석 요청에서만 포함하며 지나치게 길면 상한을 둡니다.
        prompt += "\n\n[공통 참고자료 텍스트]\n" + global_ref_text[:18000]
    payload = compile_api_payload(prompt, selected_model_name, scope="student")
    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": "교육자료 분석 전문가입니다. 원본 근거를 학생별로 정확히 분리하고 JSON 외 문장을 출력하지 마세요."},
            {"role": "user", "content": payload},
        ],
    )
    analysis = normalize_source_analysis(response.choices[0].message.content)
    inp, out, total = extract_usage_tokens(response)
    return analysis, inp, out, total


def find_student_source_summary(analysis, student_id="", student_name=""):
    if not analysis:
        return ""
    students = analysis.get("students", []) or []
    sid = (student_id or "").strip()
    name = (student_name or "").strip()
    best = None
    for item in students:
        if sid and str(item.get("id", "")).strip() == sid:
            best = item
            break
    if best is None and name:
        for item in students:
            if str(item.get("name", "")).strip() == name:
                best = item
                break
    parts = []
    if best and best.get("summary"):
        parts.append("[해당 학생 첨부자료 요약]\n" + best.get("summary", ""))
    elif analysis.get("source_summary"):
        parts.append("[첨부자료 공통 요약]\n" + analysis.get("source_summary", ""))
    if analysis.get("common_summary"):
        parts.append("[공통 참고자료 핵심]\n" + analysis.get("common_summary", ""))
    return "\n\n".join(parts)


RECORD_TYPE_DEFAULT_BYTES = {
    "교과세특": 1500,
    "동아리활동": 1500,
    "자율활동": 1500,
    "진로활동": 2100,
    "행동특성 및 종합의견": 1500,
}

def get_default_target_bytes(record_type):
    return int(RECORD_TYPE_DEFAULT_BYTES.get(record_type, 1500))

def is_byte_adjustment_request(text):
    if not text:
        return False
    lowered = str(text).lower().replace(" ", "")
    keywords = ["byte", "바이트", "분량", "줄여", "줄여줘", "늘려", "늘려줘", "맞춰", "맞추", "짧게", "길게"]
    return any(k in lowered for k in keywords)

def build_byte_adjustment_instruction(actual_bytes, target_bytes):
    actual_bytes = int(actual_bytes)
    target_bytes = int(target_bytes)
    low = int(target_bytes * 0.92)
    if actual_bytes > target_bytes:
        delta = actual_bytes - target_bytes
        return f"현재 {actual_bytes} Byte입니다. 핵심 사실은 유지하고 약 {delta} Byte 이상 줄여 최종 {low}~{target_bytes} Byte에 맞추세요."
    delta = max(0, low - actual_bytes)
    return f"현재 {actual_bytes} Byte입니다. 새로운 사실을 만들지 말고 기존 근거를 구체화하여 약 {delta} Byte 보강해 최종 {low}~{target_bytes} Byte에 맞추세요."

def generate_draft_from_summary(client, model, record_type, target_bytes, sid, name, summary, student_memo=""):
    """원본 파일 없이 캐시된 요약만으로 생기부 초안을 생성합니다."""
    system_prompt = (
        "당신은 대한민국 학교생활기록부 작성 전문가입니다. 자료에 없는 사실을 만들지 마세요. "
        "학생 이름을 주어로 시작하지 말고 곧바로 활동과 성취 특성으로 시작하세요. "
        "문장 끝은 '~함', '~임' 중심으로 작성하세요. 영어 표기와 불필요한 특수기호를 피하세요. "
        f"작성 영역은 {record_type}, 목표는 {target_bytes} Byte 이하이며 90~100%를 권장합니다. "
        "결과 생기부 문단만 출력하세요."
    )
    user_text = (
        f"[학생] {sid} {name}\n[작성 영역] {record_type}\n[목표] {target_bytes} Byte\n"
        f"[교사 메모] {student_memo or '없음'}\n\n{summary or '[첨부자료 요약 없음]'}"
    )
    response = client.chat.completions.create(
        model=model,
        messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": user_text}],
    )
    draft = response.choices[0].message.content.strip()
    inp, out, total = extract_usage_tokens(response)
    return draft, inp, out, total


def enforce_target_bytes_with_llm(client, model, text, target_bytes, record_type, max_attempts=1):
    total_in = total_out = total_tokens = 0
    revised = text
    for _ in range(max_attempts):
        _, actual = calculate_neis_bytes(revised)
        if int(target_bytes * 0.88) <= actual <= target_bytes:
            break
        adjustment_instruction = build_byte_adjustment_instruction(actual, target_bytes)
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": (
                    "학교생활기록부 문장의 길이만 조정합니다. 기존 사실과 의미를 유지하고 새 사실을 만들지 마세요. "
                    "결과 문단만 출력하세요. 바이트 계산은 프로그램이 하므로 숫자를 직접 계산하려 하지 말고, 제시된 증감량을 기준으로 문장을 조정하세요."
                )},
                {"role": "user", "content": (
                    f"작성 영역: {record_type}\n{adjustment_instruction}\n\n[현재 초안]\n{revised}"
                )}
            ]
        )
        revised = response.choices[0].message.content.strip()
        rin, rout, rtotal = extract_usage_tokens(response)
        total_in += rin
        total_out += rout
        total_tokens += rtotal
    return revised, total_in, total_out, total_tokens

FORBIDDEN_WORDS = [
    "대회", "수상", "상장", "올림피아드", "경시", "금상", "은상", "동상", "대상",
    "토익", "TOEIC", "토플", "TOEFL", "텝스", "TEPS", "HSK", "JLPT",
    "대학교", "영재교육원", "교육청", "교외", "학원", "자격증", "인증", "논문"
]

def check_forbidden_words(text):
    if not text:
        return []
    found = []
    for word in FORBIDDEN_WORDS:
        if word.lower() in text.lower():
            found.append(word)
    return list(set(found))

FORBIDDEN_WORD_GUIDE = {
    "대회": "교외 대회 참여·실적은 학교생활기록부 기재 제한 대상이 될 수 있습니다.",
    "수상": "수상 실적을 직접적으로 강조하는 표현은 기재 제한 대상이 될 수 있습니다.",
    "상장": "상장이나 수상 결과 자체를 기록하기보다 수업·활동 과정 중심으로 표현하는 것이 안전합니다.",
    "올림피아드": "교외 경시·올림피아드 관련 실적은 기재 제한 여부를 확인해야 합니다.",
    "경시": "경시대회 실적을 직접 기록하는 표현은 피하는 것이 좋습니다.",
    "금상": "수상 등급을 직접 기록하는 표현은 피하는 것이 좋습니다.",
    "은상": "수상 등급을 직접 기록하는 표현은 피하는 것이 좋습니다.",
    "동상": "수상 등급을 직접 기록하는 표현은 피하는 것이 좋습니다.",
    "대상": "문맥상 수상 등급을 뜻한다면 기재 제한에 해당할 수 있으므로 확인이 필요합니다.",
    "토익": "공인어학성적은 학교생활기록부 기재 제한 대상입니다.",
    "TOEIC": "공인어학성적은 학교생활기록부 기재 제한 대상입니다.",
    "토플": "공인어학성적은 학교생활기록부 기재 제한 대상입니다.",
    "TOEFL": "공인어학성적은 학교생활기록부 기재 제한 대상입니다.",
    "텝스": "공인어학성적은 학교생활기록부 기재 제한 대상입니다.",
    "TEPS": "공인어학성적은 학교생활기록부 기재 제한 대상입니다.",
    "HSK": "공인어학성적은 학교생활기록부 기재 제한 대상입니다.",
    "JLPT": "공인어학성적은 학교생활기록부 기재 제한 대상입니다.",
    "대학교": "특정 외부 대학·기관명을 직접 기재하는 표현은 지침상 제한될 수 있어 일반화가 필요합니다.",
    "영재교육원": "교외 기관 활동을 직접 기재하는 경우 지침 위반 가능성이 있어 확인이 필요합니다.",
    "교육청": "외부 기관명 직접 기재 여부를 확인하고 필요하면 일반화하는 것이 좋습니다.",
    "교외": "교외 활동을 직접 기재하는 경우 학교생활기록부 기재 가능 여부를 확인해야 합니다.",
    "학원": "사교육 기관 및 사교육 관련 내용은 기재하지 않는 것이 원칙입니다.",
    "자격증": "자격증 취득 실적은 기재 가능 범위를 확인해야 합니다.",
    "인증": "외부 인증·인증 실적을 의미한다면 기재 제한 여부를 확인해야 합니다.",
    "논문": "학생의 논문 게재·학술 실적을 직접 강조하는 표현은 지침상 제한될 수 있습니다.",
}

def build_validation_issue_details(text, current_bytes, target_bytes, forbidden_found):
    details = []
    if current_bytes > target_bytes:
        diff = current_bytes - target_bytes
        details.append((
            "분량 초과",
            f"현재 {current_bytes:,} Byte로 목표 {target_bytes:,} Byte보다 {diff:,} Byte 많습니다.",
            f"핵심 활동과 역량은 유지하되 반복 표현·수식어를 줄여 약 {diff:,} Byte 이상 축소해 주세요."
        ))
    elif current_bytes < int(target_bytes * 0.8):
        need = max(0, int(target_bytes * 0.9) - current_bytes)
        details.append((
            "분량 부족",
            f"현재 {current_bytes:,} Byte로 목표 {target_bytes:,} Byte의 80% 미만입니다.",
            f"새로운 사실을 만들지 말고 수행 과정·분석 근거·교사 관찰 내용을 약 {need:,} Byte 정도 보강해 주세요."
        ))
    for word in forbidden_found:
        reason = FORBIDDEN_WORD_GUIDE.get(word, "학교생활기록부 기재 제한 또는 오해 가능성이 있는 표현입니다.")
        details.append((
            f"기재 주의 표현: {word}",
            reason,
            "해당 단어를 단순 삭제하기보다 실제 수업·탐구 과정이나 학생의 행동·역량 중심 표현으로 바꾸는 것이 좋습니다."
        ))
    return details

# ==========================================
# Secrets 파일 백그라운드 파싱
# ==========================================
def load_toml_secrets():
    secrets_data = {"ACCESS_CODE": "", "DEMO_ACCESS_CODE": "", "OPENROUTER_API_KEY": ""}
    
    possible_paths = [
        os.path.join(os.getcwd(), ".streamlit", "secrets.toml"),
        os.path.join(os.path.dirname(__file__), ".streamlit", "secrets.toml"),
    ]
    
    for path in possible_paths:
        if os.path.exists(path):
            try:
                with open(path, "r", encoding="utf-8") as f:
                    content = f.read()
                    
                    code_match = re.search(r'ACCESS_CODE\s*=\s*["\']([^"\']+)["\']', content)
                    if code_match:
                        secrets_data["ACCESS_CODE"] = code_match.group(1).strip()

                    demo_code_match = re.search(r'DEMO_ACCESS_CODE\s*=\s*["\']([^"\']+)["\']', content)
                    if demo_code_match:
                        secrets_data["DEMO_ACCESS_CODE"] = demo_code_match.group(1).strip()
                        
                    key_match = re.search(r'OPENROUTER_API_KEY\s*=\s*["\']([^"\']+)["\']', content)
                    if key_match:
                        secrets_data["OPENROUTER_API_KEY"] = key_match.group(1).strip()
                break
            except Exception:
                pass
                
    try:
        if not secrets_data["OPENROUTER_API_KEY"]:
            secrets_data["OPENROUTER_API_KEY"] = str(st.secrets.get("OPENROUTER_API_KEY", "")).strip()
        if not secrets_data["ACCESS_CODE"]:
            secrets_data["ACCESS_CODE"] = str(st.secrets.get("ACCESS_CODE", "")).strip()
        if not secrets_data["DEMO_ACCESS_CODE"]:
            secrets_data["DEMO_ACCESS_CODE"] = str(st.secrets.get("DEMO_ACCESS_CODE", "")).strip()
    except Exception:
        pass
        
    return secrets_data

SECRETS = load_toml_secrets()

def is_demo_mode():
    return st.session_state.get("access_mode") == "demo"

def get_api_unavailable_message():
    if is_demo_mode():
        return "체험 모드에서는 AI/API 실행 기능이 비활성화되어 있습니다. 화면 구성과 입력 기능만 확인할 수 있습니다."
    return "API 키 설정이 완료되지 않았습니다. (.streamlit/secrets.toml 확인 필요)"

def get_openrouter_client():
    # 공유용 체험 코드는 OpenRouter API를 절대 호출하지 않습니다.
    if is_demo_mode():
        return None
    api_key = SECRETS["OPENROUTER_API_KEY"]
    if not api_key:
        return None
    return OpenAI(
        base_url="https://openrouter.ai/api/v1",
        api_key=api_key,
    )

import hashlib

# 이미지 데이터의 MD5 해시값 계산을 통한 중복 방지 헬퍼 함수
def calculate_image_hash(base64_data):
    if not base64_data:
        return ""
    payload = base64_data.split(",")[-1]
    return hashlib.md5(payload.encode("utf-8")).hexdigest()

# 모델 라인업 설정 (기능 지원 여부 및 역할 일원화 관리)
MODEL_MAP = {
    "GPT-5.6 Luna": {
        "id": "openai/gpt-5.6-luna",
        "supports_images": True,
        "role": "빠른 검토 · 가성비"
    },
    "Gemini 3.6 Flash": {
        "id": "google/gemini-3.6-flash",
        "supports_images": True,
        "role": "빠른 초안 작성"
    },
    "Claude Sonnet 5": {
        "id": "anthropic/claude-sonnet-5",
        "supports_images": True,
        "role": "생기부 작성 · 검수 권장"
    },
    "Claude Opus 5": {
        "id": "anthropic/claude-opus-5",
        "supports_images": True,
        "role": "정밀 분석"
    },
    "GPT-5.6 Sol": {
        "id": "openai/gpt-5.6-sol",
        "supports_images": True,
        "role": "고급 추론 · 복잡한 업무"
    }
}

# OpenRouter 모델 활성/사용가능 상태 검증 함수
def check_openrouter_model_availability(model_id):
    # 체험 모드는 네트워크/API 확인 자체를 생략합니다.
    if is_demo_mode():
        return True
    if "openrouter_models_cache" not in st.session_state:
        st.session_state.openrouter_models_cache = None
        
    if st.session_state.openrouter_models_cache is None:
        try:
            import urllib.request
            import json
            url = "https://openrouter.ai/api/v1/models"
            req = urllib.request.Request(
                url, 
                headers={"User-Agent": "Mozilla/5.0"}
            )
            with urllib.request.urlopen(req, timeout=4) as response:
                res_data = json.loads(response.read().decode())
                available_ids = {m.get("id") for m in res_data.get("data", []) if m.get("id")}
                st.session_state.openrouter_models_cache = available_ids
        except Exception:
            # 네트워크 타임아웃 등의 이슈로 조회가 실패하면 중단 방지를 위해 통과시킴
            return True
            
    if st.session_state.openrouter_models_cache:
        return model_id in st.session_state.openrouter_models_cache
    return True

# 2. 세션 상태 초기화
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False
if "access_mode" not in st.session_state:
    st.session_state.access_mode = None

if "user_has_manually_chosen_model" not in st.session_state:
    st.session_state.user_has_manually_chosen_model = False

if "selected_model_name" not in st.session_state:
    st.session_state.selected_model_name = "GPT-5.6 Luna"

# 파일 업로더 초기화를 위한 카운터 Key
if "uploader_key_chat" not in st.session_state:
    st.session_state.uploader_key_chat = 0
if "uploader_key_std" not in st.session_state:
    st.session_state.uploader_key_std = 0
if "uploader_key_eval" not in st.session_state:
    st.session_state.uploader_key_eval = 0

# 세션 관리 (일반 챗봇)
if "chat_sessions" not in st.session_state:
    st.session_state.chat_sessions = []
if "current_chat_idx" not in st.session_state:
    st.session_state.current_chat_idx = None

# 학생 세션 및 입력 위젯 상태 관리 (생기부)
if "student_records" not in st.session_state:
    st.session_state.student_records = []
if "current_student_idx" not in st.session_state:
    st.session_state.current_student_idx = None
if "is_new_student" not in st.session_state:
    st.session_state.is_new_student = True

# 그림자 백업 세션 (모드 탭 이동 시 폼 텍스트 보존용)
if "draft_student_id" not in st.session_state:
    st.session_state.draft_student_id = ""
if "draft_student_record_type" not in st.session_state:
    st.session_state.draft_student_record_type = "교과세특"
if "draft_student_memo" not in st.session_state:
    st.session_state.draft_student_memo = ""

if "student_id_widget" not in st.session_state:
    st.session_state.student_id_widget = st.session_state.draft_student_id
if "student_record_type_widget" not in st.session_state:
    st.session_state.student_record_type_widget = st.session_state.draft_student_record_type
if "student_memo_widget" not in st.session_state:
    st.session_state.student_memo_widget = st.session_state.draft_student_memo
if "draft_student_target_bytes" not in st.session_state:
    st.session_state.draft_student_target_bytes = 1500
if "draft_student_name" not in st.session_state:
    st.session_state.draft_student_name = ""
if "student_target_bytes_widget" not in st.session_state:
    st.session_state.student_target_bytes_widget = st.session_state.draft_student_target_bytes

# v13: v12 이전 세션에서 number_input 최소값(300)이 남아 새 교과세특에도 이어지던 상태를 1회 정리합니다.
if st.session_state.get("_student_target_schema_version", 0) < 14:
    if st.session_state.get("is_new_student", True) and st.session_state.get("current_student_idx") is None:
        current_type = st.session_state.get("draft_student_record_type", "교과세특") or "교과세특"
        default_bytes = get_default_target_bytes(current_type) if "get_default_target_bytes" in globals() else (2100 if current_type == "진로활동" else 1500)
        st.session_state.draft_student_target_bytes = default_bytes
        st.session_state.student_target_bytes_widget = default_bytes
    st.session_state._student_target_schema_version = 14

# 검수 세션 및 입력 위젯 상태 관리 (생기부 검수/진단)
if "eval_records" not in st.session_state:
    st.session_state.eval_records = []
if "current_eval_idx" not in st.session_state:
    st.session_state.current_eval_idx = None
if "is_new_eval" not in st.session_state:
    st.session_state.is_new_eval = True

if "draft_eval_text" not in st.session_state:
    st.session_state.draft_eval_text = ""
if "eval_text_widget" not in st.session_state:
    st.session_state.eval_text_widget = st.session_state.draft_eval_text

# 실행 플래그 및 초기화 확인 상태
if "is_generating_draft" not in st.session_state:
    st.session_state.is_generating_draft = False
if "is_refining_draft" not in st.session_state:
    st.session_state.is_refining_draft = False
if "is_running_eval" not in st.session_state:
    st.session_state.is_running_eval = False
if "is_generating_batch_students" not in st.session_state:
    st.session_state.is_generating_batch_students = False
if "is_analyzing_student_sources" not in st.session_state:
    st.session_state.is_analyzing_student_sources = False
if "student_source_analysis_cache" not in st.session_state:
    st.session_state.student_source_analysis_cache = {}
if "batch_failed_students" not in st.session_state:
    st.session_state.batch_failed_students = []

# API 실행은 한 번만 소비하는 one-shot 플래그로 관리합니다.
# Streamlit 상단 Stop으로 실행이 중단되더라도 다음 rerun에서 같은 요청이 자동 재실행되지 않게 합니다.
for _flag in ["_run_draft_once", "_run_refine_once", "_run_eval_once", "_run_batch_once", "_run_source_analysis_once"]:
    if _flag not in st.session_state:
        st.session_state[_flag] = False

# 이전 실행이 Stop 등으로 중단되어 busy 상태만 남았다면 다음 화면 진입 시 복구합니다.
if st.session_state.is_generating_draft and not st.session_state._run_draft_once:
    st.session_state.is_generating_draft = False
if st.session_state.is_refining_draft and not st.session_state._run_refine_once:
    st.session_state.is_refining_draft = False
if st.session_state.is_running_eval and not st.session_state._run_eval_once:
    st.session_state.is_running_eval = False
if st.session_state.is_generating_batch_students and not st.session_state._run_batch_once:
    st.session_state.is_generating_batch_students = False
if st.session_state.is_analyzing_student_sources and not st.session_state._run_source_analysis_once:
    st.session_state.is_analyzing_student_sources = False

if "confirm_new_student_reset" not in st.session_state:
    st.session_state.confirm_new_student_reset = False
if "confirm_new_eval_reset" not in st.session_state:
    st.session_state.confirm_new_eval_reset = False

# 헬퍼 함수: 새 학생 작성 전용 폼 완전 초기화 (임시 신규 항목 활성화)
def reset_student_draft_form():
    st.session_state.is_new_student = True
    st.session_state.current_student_idx = None
    st.session_state.draft_student_id = ""
    st.session_state.draft_student_name = ""
    st.session_state.draft_student_record_type = "교과세특"
    st.session_state.draft_student_memo = ""
    st.session_state.student_id_widget = ""
    st.session_state.student_record_type_widget = "교과세특"
    st.session_state.student_record_type_segmented = "교과세특"
    st.session_state.student_memo_widget = ""
    st.session_state.draft_student_target_bytes = 1500
    st.session_state.student_target_bytes_widget = 1500
    st.session_state._new_student_default_bytes_applied = "교과세특"
    st.session_state.attachments["student"] = []
    st.session_state.confirm_new_student_reset = False
    st.session_state.batch_failed_students = []
    st.session_state.pop("_pending_action_after_source_analysis", None)
    std_key = f"uploader_std_{st.session_state.uploader_key_std}"
    if std_key in st.session_state:
        del st.session_state[std_key]
    st.session_state.uploader_key_std += 1

# 헬퍼 함수: 새 검토 전용 폼 완전 초기화 (임시 신규 항목 활성화)
def reset_eval_form():
    st.session_state.is_new_eval = True
    st.session_state.current_eval_idx = None
    st.session_state.eval_result = ""
    st.session_state.eval_target_text = ""
    st.session_state.draft_eval_text = ""
    st.session_state.eval_text_widget = ""
    st.session_state.attachments["eval"] = []
    st.session_state.confirm_new_eval_reset = False
    eval_key = f"uploader_eval_{st.session_state.uploader_key_eval}"
    if eval_key in st.session_state:
        del st.session_state[eval_key]
    st.session_state.uploader_key_eval += 1

# 누적 사용 토큰 추적
if "total_tokens_used" not in st.session_state:
    st.session_state.total_tokens_used = 0

# 검수 결과 캐시 세션 (기존 단일 세션 호환 유지)
if "eval_result" not in st.session_state:
    st.session_state.eval_result = ""
if "eval_target_text" not in st.session_state:
    st.session_state.eval_target_text = ""

# 모드별 독립된 첨부파일 큐 초기화 (기존 list 세션 방어 및 마이그레이션 적용)
if "attachments" not in st.session_state or not isinstance(st.session_state.attachments, dict):
    st.session_state.attachments = {"chat": [], "student": [], "eval": []}
else:
    # 각 scope key가 누락된 경우 자동 보완
    st.session_state.attachments.setdefault("chat", [])
    st.session_state.attachments.setdefault("student", [])
    st.session_state.attachments.setdefault("eval", [])

def calculate_bytes_hash(file_bytes):
    return hashlib.md5(file_bytes).hexdigest()

def add_attachment(name, mime_type, data_base64, size_bytes, source="upload", file_type="image", scope="chat", hidden=False, parent_hash=None):
    img_hash = calculate_image_hash(data_base64)
    exists = any(item.get("hash") == img_hash for item in st.session_state.attachments[scope])
    if not exists:
        st.session_state.attachments[scope].append({
            "name": name,
            "mime_type": mime_type,
            "data": data_base64,
            "size": size_bytes,
            "size_kb": size_bytes / 1024.0,
            "source": source,
            "type": file_type,
            "hash": img_hash,
            "hidden": bool(hidden),
            "parent_hash": parent_hash
        })
        return True
    return False

def add_document_attachment(name, mime_type, text_content, size_bytes, file_hash, scope="chat"):
    exists = any(item.get("hash") == file_hash for item in st.session_state.attachments[scope])
    if not exists:
        st.session_state.attachments[scope].append({
            "name": name,
            "mime_type": mime_type,
            "data": text_content,
            "size": size_bytes,
            "size_kb": size_bytes / 1024.0,
            "source": "upload",
            "type": "document",
            "hash": file_hash
        })
        return True
    return False

PDF_IMAGE_MARKER = "[이미지형 PDF]"
PDF_VISION_MAX_PAGES = None  # v15: 이미지형 PDF는 모든 페이지 분석
PDF_VISION_MAX_WIDTH = 768
PDF_VISION_JPEG_QUALITY = 62

@st.cache_data(ttl="2h", max_entries=20, show_spinner=False)
def render_pdf_pages_as_images(file_bytes, max_pages=PDF_VISION_MAX_PAGES):
    """텍스트 레이어가 없는 PDF의 모든 페이지를 Vision 입력용 경량 JPEG data URL 목록으로 변환합니다."""
    import base64
    from PIL import Image
    rendered = []
    total_pages = 0

    # 1순위: PyMuPDF가 있으면 빠르게 렌더링
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        total_pages = len(doc)
        page_count = total_pages if max_pages is None else min(total_pages, max_pages)
        for page_idx in range(page_count):
            page = doc.load_page(page_idx)
            pix = page.get_pixmap(matrix=fitz.Matrix(1.0, 1.0), alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            if img.width > PDF_VISION_MAX_WIDTH:
                ratio = PDF_VISION_MAX_WIDTH / float(img.width)
                img = img.resize((PDF_VISION_MAX_WIDTH, max(1, int(img.height * ratio))))
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=PDF_VISION_JPEG_QUALITY, optimize=True)
            raw = buf.getvalue()
            data_url = "data:image/jpeg;base64," + base64.b64encode(raw).decode("utf-8")
            rendered.append((page_idx + 1, data_url, len(raw)))
        doc.close()
        return rendered, total_pages
    except Exception:
        rendered = []

    # 2순위: 별도 PyMuPDF 설치가 없는 환경에서는 기존 pdfplumber 렌더러 사용
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            total_pages = len(pdf.pages)
            pages_to_render = pdf.pages if max_pages is None else pdf.pages[:max_pages]
            for page_idx, page in enumerate(pages_to_render):
                page_image = page.to_image(resolution=80).original.convert("RGB")
                img = page_image
                if img.width > PDF_VISION_MAX_WIDTH:
                    ratio = PDF_VISION_MAX_WIDTH / float(img.width)
                    img = img.resize((PDF_VISION_MAX_WIDTH, max(1, int(img.height * ratio))))
                buf = io.BytesIO()
                img.save(buf, format="JPEG", quality=PDF_VISION_JPEG_QUALITY, optimize=True)
                raw = buf.getvalue()
                data_url = "data:image/jpeg;base64," + base64.b64encode(raw).decode("utf-8")
                rendered.append((page_idx + 1, data_url, len(raw)))
        return rendered, total_pages
    except Exception:
        return [], total_pages

def add_pdf_attachment_with_vision(uploaded_file, scope, file_bytes, file_hash):
    """PDF에 텍스트가 없으면 페이지를 이미지로 변환해 Vision 모델이 읽을 수 있게 첨부합니다."""
    text_content = parse_uploaded_file(uploaded_file)
    add_document_attachment(
        name=uploaded_file.name,
        mime_type=uploaded_file.type,
        text_content=text_content,
        size_bytes=uploaded_file.size,
        file_hash=file_hash,
        scope=scope
    )
    if text_content.startswith(PDF_IMAGE_MARKER):
        pages, _ = render_pdf_pages_as_images(file_bytes)
        for page_no, data_url, image_size in pages:
            add_attachment(
                name=f"{uploaded_file.name} · {page_no}쪽",
                mime_type="image/jpeg",
                data_base64=data_url,
                size_bytes=image_size,
                source="pdf_page",
                file_type="image",
                scope=scope,
                hidden=True,
                parent_hash=file_hash
            )
        return len(pages)
    return 0

def delete_attachment_group(scope, item_hash):
    st.session_state.attachments[scope] = [
        item for item in st.session_state.attachments[scope]
        if item.get("hash") != item_hash and item.get("parent_hash") != item_hash
    ]

def render_attachments_panel(uploader_key, scope="chat"):
    paste_key = f"paste_listener_{scope}"
    
    # 1. paste listener callback (scope 인식)
    def handle_pasted_image():
        state = st.session_state.get(paste_key)
        if state and hasattr(state, "pasted_image") and state.pasted_image:
            img_data = state.pasted_image
            add_attachment(
                name=img_data["name"],
                mime_type=img_data["type"],
                data_base64=img_data["data"],
                size_bytes=img_data["size"],
                source="paste",
                file_type="image",
                scope=scope
            )
            
    paste_listener(key=paste_key, on_pasted_image=handle_pasted_image)
    
    current_list = st.session_state.attachments[scope]
    
    # 2. Scope별 구분 라벨
    if scope == "student":
        st.caption("현재 학생 첨부자료 (파일 선택 또는 Ctrl+V 캡처)")
    elif scope == "eval":
        st.caption("검토 첨부자료 (파일 선택 또는 Ctrl+V 캡처)")

    # 3. 항상 노출형 콤팩트 파일 업로더
    uploaded_files = st.file_uploader(
        "파일 선택 또는 Ctrl+V로 이미지 붙여넣기",
        accept_multiple_files=True,
        key=uploader_key,
        label_visibility="collapsed"
    )
    
    if uploaded_files:
        for f in uploaded_files:
            file_bytes = f.getvalue()
            file_hash = calculate_bytes_hash(file_bytes)
            
            if f.name.lower().endswith(('.png', '.jpg', '.jpeg', '.webp')):
                import base64
                b64_data = f"data:{f.type};base64," + base64.b64encode(file_bytes).decode("utf-8")
                add_attachment(
                    name=f.name,
                    mime_type=f.type,
                    data_base64=b64_data,
                    size_bytes=f.size,
                    source="upload",
                    file_type="image",
                    scope=scope
                )
            elif f.name.lower().endswith('.pdf'):
                page_images = add_pdf_attachment_with_vision(f, scope, file_bytes, file_hash)
                if page_images:
                    st.caption(f"이미지형 PDF 감지: 최대 {page_images}쪽을 768px JPEG로 경량화해 최초 1회 분석합니다.")
            else:
                text_content = parse_uploaded_file(f)
                add_document_attachment(
                    name=f.name,
                    mime_type=f.type,
                    text_content=text_content,
                    size_bytes=f.size,
                    file_hash=file_hash,
                    scope=scope
                )

    # 4. 첨부자료가 존재할 때만 콤팩트 칩스 바 렌더링
    visible_items = [(idx, item) for idx, item in enumerate(current_list) if not item.get("hidden", False)]
    if visible_items:
        max_cols = min(len(visible_items), 6)
        cols = st.columns(max_cols)
        for visible_idx, (original_idx, item) in enumerate(visible_items):
            with cols[visible_idx % max_cols]:
                if item["type"] == "image":
                    st.image(item["data"], width=45)
                    btn_label = f"삭제 {item['name'][:8]}"
                else:
                    btn_label = f"삭제 [문서]{item['name'][:8]}"
                if st.button(btn_label, key=f"del_att_{uploader_key}_{original_idx}", help=f"{item['name']} ({item['size_kb']:.0f}KB) 삭제"):
                    delete_attachment_group(scope, item.get("hash"))
                    st.rerun()

def compile_api_payload(prompt, selected_model_name, scope="chat"):
    images = []
    text_contexts = []
    
    for item in st.session_state.attachments[scope]:
        if item["type"] == "image":
            images.append(item)
        elif item["type"] == "document":
            text_contexts.append(f"\n\n[첨부문서: {item['name']}]\n" + item["data"])
            
    model_supports_img = MODEL_MAP[selected_model_name].get("supports_images", False)
    
    if images and not model_supports_img:
        st.error(f"현재 선택한 모델 ({selected_model_name})은 이미지 분석을 지원하지 않습니다. 이미지를 제외하거나 다른 모델을 선택해 주세요.")
        st.stop()
        
    doc_text_combined = "".join(text_contexts)
    final_text = prompt + doc_text_combined if doc_text_combined else prompt
    
    content_payload = []
    content_payload.append({"type": "text", "text": final_text})
    
    for img in images:
        clean_b64 = img["data"].split(",")[-1]
        content_payload.append({
            "type": "image_url",
            "image_url": {
                "url": f"data:{img['mime_type']};base64,{clean_b64}"
            }
        })
        
    return content_payload

# 일반 챗봇 대화 문맥 제한: 최근 10개 메시지만 API에 재전송
# (약 5턴 분량이며, 과거 첨부 원문/이미지는 반복 전송하지 않음)
CHAT_HISTORY_MAX_MESSAGES = 10

def extract_usage_tokens(response):
    """OpenRouter/OpenAI 호환 usage에서 입력/출력/합계 토큰을 안전하게 추출합니다."""
    usage = getattr(response, "usage", None)
    if not usage:
        return 0, 0, 0

    def _get(name, fallback=0):
        if isinstance(usage, dict):
            value = usage.get(name, fallback)
        else:
            value = getattr(usage, name, fallback)
        try:
            return int(value or 0)
        except (TypeError, ValueError):
            return 0

    input_tokens = _get("prompt_tokens") or _get("input_tokens")
    output_tokens = _get("completion_tokens") or _get("output_tokens")
    total_tokens = _get("total_tokens")
    if total_tokens <= 0:
        total_tokens = input_tokens + output_tokens
    return input_tokens, output_tokens, total_tokens

def compact_history_content(message):
    """과거 사용자 메시지에서는 첨부 원문/Base64를 제거하고 질문 요약만 문맥으로 사용합니다."""
    if message.get("history_text") is not None:
        return message.get("history_text") or ""

    content = message.get("content", "")
    if isinstance(content, str):
        return content

    if isinstance(content, list):
        text_parts = []
        had_image = False
        for part in content:
            if not isinstance(part, dict):
                continue
            if part.get("type") == "text":
                txt = part.get("text", "")
                # 기존 세션 호환: 문서 원문이 합쳐진 경우 첨부문서 시작 전까지만 유지
                if "\n\n[첨부문서:" in txt:
                    txt = txt.split("\n\n[첨부문서:", 1)[0]
                if txt.strip():
                    text_parts.append(txt.strip())
            elif part.get("type") == "image_url":
                had_image = True
        result = "\n".join(text_parts).strip()
        if had_image:
            result += ("\n" if result else "") + "[이전 질문에 이미지 첨부자료가 있었음 - 원본은 재전송하지 않음]"
        return result

    return str(content)

# ==========================================
# v16 일반 챗봇: 파일 메타데이터 · 분석 캐시 · Excel 실파일 출력
# ==========================================
CHAT_FILE_CACHE_MAX_CHARS = 30000


def wants_excel_output(prompt):
    text = (prompt or "").lower()
    keywords = ["엑셀", "excel", ".xlsx", "xlsx", "스프레드시트", "spreadsheet", "엑셀파일", "엑셀 파일", "엑셀로 저장", "엑셀로 정리"]
    return any(k in text for k in keywords)


def get_pdf_page_count(file_bytes):
    try:
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        count = len(doc)
        doc.close()
        return int(count)
    except Exception:
        try:
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                return len(pdf.pages)
        except Exception:
            return 0


def build_file_metadata_context(file_meta):
    if not file_meta:
        return ""
    lines = ["[앱이 직접 확인한 현재 첨부파일 정보]"]
    for meta in file_meta:
        name = meta.get("name", "첨부파일")
        kind = meta.get("kind", "file")
        if kind == "pdf":
            total = int(meta.get("total_pages", 0) or 0)
            sent = int(meta.get("sent_pages", total) or total)
            lines.append(f"- {name}: PDF 총 {total}페이지, 현재 요청에 {sent}/{total}페이지 전달")
        elif kind == "image":
            lines.append(f"- {name}: 이미지 1장 전달")
        else:
            lines.append(f"- {name}: {meta.get('description', '문서 본문을 추출하여 전달')}")
    lines.append("위 페이지 수와 전달 수는 앱이 직접 확인한 값이므로 임의로 다른 숫자를 추정하지 마세요.")
    return "\n".join(lines)


def build_chat_system_prompt(current_chat, current_message):
    file_meta = current_message.get("file_meta", []) if isinstance(current_message, dict) else []
    meta_text = build_file_metadata_context(file_meta)
    cached_source = (current_chat.get("source_cache_text") or "").strip()
    excel_intent = bool(current_message.get("excel_intent")) if isinstance(current_message, dict) else False
    rules = (
        "당신은 교사를 보조하는 AI 에이전트 'Chat PSDongSung'입니다.\n"
        "첨부자료 처리 원칙:\n"
        "1. 현재 요청에 첨부자료가 실제로 제공되었다면 먼저 자료를 직접 확인하고 답하세요. 충분히 전달된 파일을 보지 못했다고 말하거나 재첨부를 요구하지 마세요.\n"
        "2. PDF의 총 페이지 수와 전달 페이지 수는 아래 앱 메타데이터를 최우선 사실로 사용하세요. 페이지 수를 추측하지 마세요.\n"
        "3. 이미지형 PDF와 사진은 보이는 글자, 표, 손글씨, 캡션을 가능한 한 적극적으로 판독하세요. 불필요하게 '판독불가'를 남발하지 마세요.\n"
        "4. 이름·학번·진로·수치처럼 식별 가능한 정보는 자료가 뒷받침하는 범위에서 사용하고, 근거 없이 새 사실을 만들지는 마세요.\n"
        "5. 과거 첨부 원본이 이번 요청에 재전송되지 않은 경우에는 아래 파일 분석 캐시를 이전 자료의 근거로 사용하세요. 캐시에 없는 세부 내용을 본 것처럼 꾸미지 마세요.\n"
        "6. 사용자가 파일 변환·표 정리·엑셀 저장을 요청하면 설명만 하지 말고 요청한 결과물을 완성하는 데 집중하세요."
    )
    parts = [rules]
    if meta_text:
        parts.append(meta_text)
    if cached_source:
        parts.append("[이 대화에서 이전 첨부자료를 1회 분석해 저장한 파일 분석 캐시]\n" + cached_source[:CHAT_FILE_CACHE_MAX_CHARS])
    if excel_intent:
        parts.append(
            "[Excel 출력 모드 - 매우 중요]\n"
            "사용자가 실제 Excel 파일을 원합니다. 응답 전체를 반드시 아래 태그 하나로만 출력하세요. 태그 밖에는 어떤 설명이나 마크다운도 쓰지 마세요.\n"
            '<EXCEL_DATA>{"filename":"정리결과.xlsx","sheets":[{"name":"정리결과","columns":["열1","열2"],"rows":[["값1","값2"]]}]}</EXCEL_DATA>\n'
            "- JSON은 반드시 유효해야 합니다.\n"
            "- columns에는 사용자의 요청과 자료 구조에 맞는 실제 열 이름을 넣으세요.\n"
            "- rows에는 첨부자료에서 확인한 데이터를 빠짐없이 넣으세요.\n"
            "- 여러 학생이면 학생 1명당 1행을 우선합니다.\n"
            "- 원문 전사가 요청된 경우 임의 요약보다 원문 표현을 우선합니다.\n"
            "- 여러 시트가 유용하면 sheets 배열에 추가해도 됩니다."
        )
    elif file_meta:
        parts.append(
            "[후속 질문용 파일 분석 캐시 생성]\n"
            "현재 첨부자료를 읽은 뒤 일반 답변 마지막에 아래 태그를 추가하세요.\n"
            '<FILE_CACHE>{"summary":"후속 질문에서 원본 이미지를 다시 보내지 않아도 될 정도로 페이지/학생/항목별 핵심 사실과 원문 표현을 정보 밀도 높게 정리"}</FILE_CACHE>\n'
            "캐시는 최대한 구체적으로 작성하되 같은 내용을 반복하지 마세요."
        )
    return "\n\n".join(parts)


def parse_hidden_chat_payloads(raw_text):
    raw = raw_text or ""
    excel_data = None
    file_cache = None
    m = re.search(r"<EXCEL_DATA>\s*(\{.*?\})\s*</EXCEL_DATA>", raw, flags=re.DOTALL)
    if m:
        try:
            excel_data = json.loads(m.group(1))
        except Exception:
            pass
    m = re.search(r"<FILE_CACHE>\s*(\{.*?\})\s*</FILE_CACHE>", raw, flags=re.DOTALL)
    if m:
        try:
            payload = json.loads(m.group(1))
            file_cache = str(payload.get("summary", "") if isinstance(payload, dict) else payload).strip()
        except Exception:
            file_cache = m.group(1).strip()
    visible = re.sub(r"<EXCEL_DATA>.*?</EXCEL_DATA>", "", raw, flags=re.DOTALL)
    visible = re.sub(r"<FILE_CACHE>.*?</FILE_CACHE>", "", visible, flags=re.DOTALL).strip()
    return visible, excel_data, file_cache


def _safe_sheet_name(name, used):
    clean = re.sub(r"[\\/*?:\[\]]", "_", str(name or "정리결과")).strip()[:31] or "정리결과"
    base = clean
    n = 2
    while clean in used:
        suffix = f"_{n}"
        clean = base[:31-len(suffix)] + suffix
        n += 1
    used.add(clean)
    return clean


def build_excel_bytes(excel_spec):
    if not isinstance(excel_spec, dict):
        return None, None
    sheets = excel_spec.get("sheets")
    if not isinstance(sheets, list) or not sheets:
        return None, None
    filename = str(excel_spec.get("filename") or "정리결과.xlsx").strip()
    if not filename.lower().endswith(".xlsx"):
        filename += ".xlsx"
    output = io.BytesIO()
    used = set()
    try:
        with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
            workbook = writer.book
            wrap_fmt = workbook.add_format({"text_wrap": True, "valign": "top"})
            header_fmt = workbook.add_format({"bold": True, "text_wrap": True, "valign": "top"})
            wrote = False
            for sheet in sheets:
                if not isinstance(sheet, dict):
                    continue
                columns = sheet.get("columns") or []
                rows = sheet.get("rows") or []
                if not isinstance(columns, list) or not columns:
                    continue
                normalized_rows = []
                for row in rows if isinstance(rows, list) else []:
                    if isinstance(row, dict):
                        normalized_rows.append([row.get(col, "") for col in columns])
                    elif isinstance(row, list):
                        normalized_rows.append((row + [""] * len(columns))[:len(columns)])
                df = pd.DataFrame(normalized_rows, columns=[str(c) for c in columns])
                sheet_name = _safe_sheet_name(sheet.get("name"), used)
                df.to_excel(writer, index=False, sheet_name=sheet_name)
                wrote = True
                ws = writer.sheets[sheet_name]
                ws.freeze_panes(1, 0)
                if len(columns):
                    ws.autofilter(0, 0, max(0, len(df)), len(columns)-1)
                for c_idx, col in enumerate(columns):
                    vals = [str(col)] + [str(v) for v in df.iloc[:, c_idx].fillna("").tolist()]
                    width = min(45, max(10, max((len(v) for v in vals), default=10) + 2))
                    ws.set_column(c_idx, c_idx, width, wrap_fmt)
                ws.set_row(0, None, header_fmt)
            if not wrote:
                return None, None
        return output.getvalue(), filename
    except Exception:
        return None, None


def excel_spec_to_cache(excel_spec):
    if not isinstance(excel_spec, dict):
        return ""
    chunks = []
    current_len = 0
    for sheet in excel_spec.get("sheets", []) or []:
        if not isinstance(sheet, dict):
            continue
        columns = [str(c) for c in (sheet.get("columns") or [])]
        if not columns:
            continue
        head = f"[시트: {sheet.get('name', '정리결과')}]\n열: " + " | ".join(columns)
        chunks.append(head); current_len += len(head)
        for row in (sheet.get("rows") or []):
            if isinstance(row, dict):
                vals = [str(row.get(c, "")) for c in columns]
            elif isinstance(row, list):
                vals = [str(v) for v in row[:len(columns)]]
            else:
                continue
            line = " | ".join(vals)
            chunks.append(line); current_len += len(line)
            if current_len >= CHAT_FILE_CACHE_MAX_CHARS:
                break
        if current_len >= CHAT_FILE_CACHE_MAX_CHARS:
            break
    return "\n".join(chunks)[:CHAT_FILE_CACHE_MAX_CHARS]

# ==========================================
# HWPX 안전 파싱 및 첨부 텍스트 상한
# ==========================================
MAX_ATTACHMENT_TEXT_CHARS = 30000

def _xml_local_name(tag):
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag

def parse_hwpx_content(file_bytes):
    """HWPX(ZIP/XML)에서 실제 문단 텍스트만 추출한다."""
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            section_names = [
                name for name in zf.namelist()
                if re.match(r"^Contents/section\d+\.xml$", name, flags=re.IGNORECASE)
            ]
            section_names.sort(key=lambda name: int(re.search(r"section(\d+)\.xml$", name, re.IGNORECASE).group(1)))
            if not section_names:
                return "[HWPX 본문 XML을 찾지 못했습니다. 파일이 손상되었거나 지원되지 않는 형식일 수 있습니다.]"

            paragraphs = []
            for section_name in section_names:
                root = ET.fromstring(zf.read(section_name))
                for elem in root.iter():
                    if _xml_local_name(elem.tag) != "p":
                        continue
                    fragments = []
                    for child in elem.iter():
                        local = _xml_local_name(child.tag)
                        if local == "t" and child.text:
                            fragments.append(child.text)
                        elif local in {"lineBreak", "br"}:
                            fragments.append("\n")
                        elif local == "tab":
                            fragments.append("\t")
                    line = "".join(fragments).strip()
                    if line:
                        paragraphs.append(line)

            text = "\n".join(paragraphs).strip()
            if not text:
                return "[HWPX에서 읽을 수 있는 본문 텍스트를 찾지 못했습니다.]"
            return text
    except zipfile.BadZipFile:
        return "[HWPX 파일 구조를 읽지 못했습니다. 파일이 손상되었거나 실제 HWPX 형식이 아닐 수 있습니다.]"
    except Exception as e:
        return f"[HWPX 본문 파싱 실패: {str(e)}]"

def limit_attachment_text(text, file_name):
    if not text:
        return text
    if len(text) <= MAX_ATTACHMENT_TEXT_CHARS:
        return text
    return (
        text[:MAX_ATTACHMENT_TEXT_CHARS]
        + f"\n\n[첨부문서 {file_name}: 본문이 너무 길어 앞 {MAX_ATTACHMENT_TEXT_CHARS:,}자까지만 AI에 전달했습니다. "
          "필요한 부분만 별도 파일로 첨부하면 더 정확하고 비용도 줄어듭니다.]"
    )

# ==========================================
# 파일 통합 파싱 및 캐싱 최적화 (버퍼링 완벽 해결)
# ==========================================
@st.cache_data(ttl="1h", max_entries=50)
def parse_file_content(file_name, file_bytes):
    extracted_text = ""
    file_name = file_name.lower()
    try:
        file_io = io.BytesIO(file_bytes)
        
        if file_name.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(file_io).fillna("")
            # v13: 표 구조를 보존한 TSV로 저장해 다중 학생 자료를 AI 호출 없이 직접 분리할 수 있게 함
            extracted_text = "[STRUCTURED_TABLE_TSV]\n" + df.astype(str).to_csv(sep="\t", index=False)
        elif file_name.endswith('.csv'):
            df = pd.read_csv(file_io).fillna("")
            extracted_text = "[STRUCTURED_TABLE_TSV]\n" + df.astype(str).to_csv(sep="\t", index=False)
        elif file_name.endswith('.pdf'):
            page_count = 0
            with pdfplumber.open(file_io) as pdf:
                page_count = len(pdf.pages)
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        extracted_text += text + "\n"
            # 프레젠테이션/스캔 PDF처럼 텍스트 레이어가 없으면 Vision 처리 대상으로 표시
            if len(extracted_text.strip()) < 40:
                extracted_text = (
                    f"{PDF_IMAGE_MARKER} 총 {page_count}쪽. 텍스트 레이어가 없어 "
                    "각 페이지 화면 이미지를 함께 첨부했습니다. 슬라이드의 제목, 본문, 표와 시각 요소를 직접 읽어 분석하세요."
                )
        elif file_name.endswith('.docx'):
            doc = Document(file_io)
            for paragraph in doc.paragraphs:
                extracted_text += paragraph.text + "\n"
        elif file_name.endswith('.pptx'):
            try:
                prs = Presentation(file_io)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            except Exception:
                extracted_text = f"[PPTX 문서 첨부됨: {file_name}]"
        elif file_name.endswith('.hwpx'):
            extracted_text = parse_hwpx_content(file_bytes)
        elif file_name.endswith('.hwp'):
            extracted_text = (
                "[구형 HWP 파일은 바이너리 원문을 안전하게 추출하지 못해 AI에 전송하지 않았습니다. "
                "한글에서 HWPX, DOCX 또는 PDF로 저장한 뒤 다시 첨부해 주세요.]"
            )
        elif file_name.endswith(('.png', '.jpg', '.jpeg')):
            extracted_text = f"[이미지 파일 첨부됨: {file_name}]"
        else:
            extracted_text = f"[첨부 파일: {file_name}]"
    except Exception as e:
        extracted_text = f"파싱 참고: {str(e)}"
    return limit_attachment_text(extracted_text, file_name)

def parse_uploaded_file(uploaded_file):
    if uploaded_file is None:
        return ""
    
    # 텍스트 입력창 타이핑 중 랙 걸리는 현상을 방지하는 st.session_state 2차 캐싱
    # 파일명/크기가 같아도 내용이 다른 파일을 혼동하지 않도록 실제 바이트 해시를 키로 사용
    file_bytes = uploaded_file.getvalue()
    file_key = f"parsed_{calculate_bytes_hash(file_bytes)}"
    if file_key not in st.session_state:
        st.session_state[file_key] = parse_file_content(uploaded_file.name, file_bytes)
        
    return st.session_state[file_key]

# 3. 영화/드라마 보안 시스템 스타일 콤팩트 중앙 로그인 화면
if not st.session_state.authenticated:
    st.markdown("""
        <style>
        /* 중앙 컬럼을 반투명 보안 대시보드 카드로 연출 */
        div[data-testid="column"]:nth-of-type(2) {
            background: #0f172a !important;
            border: 1px solid #334155 !important;
            border-radius: 16px !important;
            padding: 2.2rem 2rem 1.8rem 2rem !important;
            box-shadow: 0 10px 25px -5px rgba(0, 0, 0, 0.4), 0 8px 10px -6px rgba(0, 0, 0, 0.3) !important;
            margin-top: 2rem !important;
        }
        </style>
    """, unsafe_allow_html=True)
    
    col_l1, col_l2, col_l3 = st.columns([1, 1.2, 1])
    with col_l2:
        st.markdown("""
            <div class="login-card">
                <div class="login-title">Chat <span class="brand-accent">PSDongSung</span></div>
                <div class="login-sub">RESTRICTED ACCESS • TEACHER ONLY</div>
            </div>
        """, unsafe_allow_html=True)
        
        access_code = st.text_input("접속 코드 입력", type="password", placeholder="보안 코드를 입력하세요", label_visibility="collapsed")
        if st.button("SYSTEM ACCESS", type="primary", use_container_width=True):
            target_code = SECRETS["ACCESS_CODE"]
            demo_code = SECRETS["DEMO_ACCESS_CODE"]
            if not target_code and not demo_code:
                st.error("ACCESS_CODE 또는 DEMO_ACCESS_CODE가 설정되지 않았습니다. .streamlit/secrets.toml을 확인해 주세요.")
            elif target_code and access_code == target_code:
                st.session_state.authenticated = True
                st.session_state.access_mode = "full"
                st.rerun()
            elif demo_code and access_code == demo_code:
                st.session_state.authenticated = True
                st.session_state.access_mode = "demo"
                st.rerun()
            else:
                st.error("접속 코드가 올바르지 않습니다.")
    st.stop()

if is_demo_mode():
    st.info("🧪 공유용 체험 모드 · AI/API 호출은 차단되어 있으며 화면 구성과 입력 기능만 확인할 수 있습니다.")

# 최초 접속 시 기본 세션 생성
if st.session_state.current_chat_idx is None and not st.session_state.chat_sessions:
    st.session_state.chat_sessions.append({"title": "새 대화", "messages": []})
    st.session_state.current_chat_idx = 0

# 4. 상단 정적 영역 (1단: flex: 0 0 auto - 로고 / 모델 선택 / 모드 탭 / 구분선)
with st.container(key="top_area"):
    mode_from_state = st.session_state.get("mode_control_widget", "일반 챗봇")
    model_keys = list(MODEL_MAP.keys())

    # 모드가 실제로 바뀌는 순간 권장 모델을 강제 적용. 이후에는 사용자가 수동 변경 가능.
    prev_mode_for_model = st.session_state.get("_last_mode_for_model")
    if prev_mode_for_model != mode_from_state:
        forced_model = "Claude Sonnet 5" if mode_from_state in ["생기부 작성", "생기부 검수/진단"] else "GPT-5.6 Luna"
        st.session_state.selected_model_name = forced_model
        st.session_state["model_selectbox_widget"] = forced_model
        st.session_state.user_has_manually_chosen_model = False
        st.session_state._last_mode_for_model = mode_from_state

    # 세션에 기록된 모델명의 인덱스 검색
    default_idx = 0
    if st.session_state.selected_model_name in model_keys:
        default_idx = model_keys.index(st.session_state.selected_model_name)

    col1, col2 = st.columns([3, 1], vertical_alignment="center")
    with col1:
        st.markdown("""
            <div class="brand-header">
                <div class="brand-title">Chat <span class="brand-accent">PSDongSung</span></div>
                <div class="brand-sub">교사 전용 스마트 AI 에이전트</div>
            </div>
        """, unsafe_allow_html=True)
    with col2:
        def on_model_change():
            st.session_state.user_has_manually_chosen_model = True

        selected_model_name = st.selectbox(
            "모델 선택",
            model_keys,
            index=default_idx,
            label_visibility="collapsed",
            key="model_selectbox_widget",
            on_change=on_model_change
        )
        st.session_state.selected_model_name = selected_model_name
        selected_model_info = MODEL_MAP[selected_model_name]
        selected_model = selected_model_info["id"]
        role_text = selected_model_info["role"]
        
        st.caption(role_text)

    # 모드 선택
    def _on_mode_change():
        new_mode = st.session_state.get("mode_control_widget", "일반 챗봇")
        # v14: 일반 챗봇 등 다른 모드에서 생기부 작성으로 처음 들어오는 순간,
        # 오래된 number_input 값(예: 300)이 화면에 먼저 나타나지 않도록 기본값을 선동기화합니다.
        if new_mode == "생기부 작성" and st.session_state.get("is_new_student", True) and st.session_state.get("current_student_idx") is None:
            current_type = st.session_state.get("draft_student_record_type", "교과세특") or "교과세특"
            if current_type not in ["교과세특", "동아리활동", "자율활동", "진로활동", "행동특성 및 종합의견"]:
                current_type = "교과세특"
            default_bytes = get_default_target_bytes(current_type)
            st.session_state.draft_student_record_type = current_type
            st.session_state.student_record_type_widget = current_type
            st.session_state.student_record_type_segmented = current_type
            st.session_state.draft_student_target_bytes = default_bytes
            st.session_state.student_target_bytes_widget = default_bytes
            st.session_state._new_student_default_bytes_applied = current_type

    mode = st.segmented_control(
        "모드",
        ["일반 챗봇", "생기부 작성", "생기부 검수/진단"],
        default="일반 챗봇",
        label_visibility="collapsed",
        key="mode_control_widget",
        on_change=_on_mode_change
    )

# 5. 사이드바 구성
with st.sidebar:
    st.markdown("""
        <div class="sidebar-brand-title">
            Chat <span class="brand-accent">PSDongSung</span>
        </div>
    """, unsafe_allow_html=True)
    
    if mode == "일반 챗봇":
        if st.button("새 대화 시작", use_container_width=True):
            chat_key = f"uploader_chat_{st.session_state.uploader_key_chat}"
            if chat_key in st.session_state:
                del st.session_state[chat_key]
            st.session_state.uploader_key_chat += 1
            st.session_state.attachments["chat"] = []
            
            new_idx = len(st.session_state.chat_sessions)
            st.session_state.chat_sessions.append({"title": "새 대화", "messages": []})
            st.session_state.current_chat_idx = new_idx
            st.rerun()
            
        st.subheader("대화 목록")
        for idx, chat in enumerate(st.session_state.chat_sessions):
            btn_label = f"{chat['title']}"
            is_active = (idx == st.session_state.current_chat_idx)
            if st.button(
                btn_label, 
                key=f"chat_btn_{idx}", 
                use_container_width=True, 
                type="primary" if is_active else "secondary"
            ):
                chat_key = f"uploader_chat_{st.session_state.uploader_key_chat}"
                if chat_key in st.session_state:
                    del st.session_state[chat_key]
                st.session_state.uploader_key_chat += 1
                st.session_state.attachments["chat"] = []
                
                st.session_state.current_chat_idx = idx
                st.rerun()
                
    elif mode == "생기부 작성":
        if st.button("새 학생 작성", use_container_width=True):
            has_unsaved_input = (st.session_state.is_new_student) and (
                bool(st.session_state.get("student_id_widget")) or bool(st.session_state.get("student_memo_widget"))
            )
            if has_unsaved_input:
                st.session_state.confirm_new_student_reset = True
            else:
                reset_student_draft_form()
            st.rerun()

        if st.session_state.get("confirm_new_student_reset", False):
            st.warning("작성 중인 내용이 있습니다. 새 작성을 시작하면 현재 입력 내용이 초기화됩니다.")
            col_c1, col_c2 = st.columns(2)
            with col_c1:
                if st.button("새 작성 계속", type="primary", key="btn_confirm_std_reset"):
                    reset_student_draft_form()
                    st.rerun()
            with col_c2:
                if st.button("취소", key="btn_cancel_std_reset"):
                    st.session_state.confirm_new_student_reset = False
                    st.rerun()

        st.subheader("학생 목록")
        
        # 1. 저장된 기존 학생 목록 (순서 및 위치 고정)
        for idx, student in enumerate(st.session_state.student_records):
            student_name_label = clean_student_name(student.get("name", ""), student.get("id_val", ""))
            btn_label = f"{student['id_val']} {student_name_label}".strip()
            is_active = (not st.session_state.is_new_student) and (idx == st.session_state.current_student_idx)
            if st.button(
                btn_label, 
                key=f"std_btn_{idx}_{student['id_val']}", 
                use_container_width=True, 
                type="primary" if is_active else "secondary"
            ):
                st.session_state.is_new_student = False
                st.session_state.current_student_idx = idx
                st.session_state.draft_student_id = student["id_val"]
                st.session_state.draft_student_name = student.get("name", "")
                loaded_record_type = student.get("record_type", "교과세특")
                if " (" in loaded_record_type:
                    loaded_record_type = loaded_record_type.split(" (")[0]
                st.session_state.draft_student_record_type = loaded_record_type
                st.session_state.draft_student_memo = student["memo"]
                st.session_state.draft_student_target_bytes = int(student.get("target_bytes", 1500))
                st.session_state.student_id_widget = student["id_val"]
                st.session_state.student_record_type_widget = loaded_record_type
                st.session_state.student_record_type_segmented = loaded_record_type
                st.session_state.student_target_bytes_widget = int(student.get("target_bytes", 1500))
                st.session_state.student_memo_widget = student["memo"]
                st.session_state.confirm_new_student_reset = False
                std_key = f"uploader_std_{st.session_state.uploader_key_std}"
                if std_key in st.session_state:
                    del st.session_state[std_key]
                st.session_state.uploader_key_std += 1
                st.session_state.attachments["student"] = []
                st.rerun()

        # 2. 임시 신규 학생 항목 (목록 맨 아래에 파란색 활성화 표시)
        if st.session_state.is_new_student:
            st.button("새 학생", key="temp_new_student_btn", use_container_width=True, type="primary")

    elif mode == "생기부 검수/진단":
        if st.button("새 검토", use_container_width=True):
            has_unsaved_eval = (st.session_state.is_new_eval) and bool(st.session_state.get("eval_text_widget"))
            if has_unsaved_eval:
                st.session_state.confirm_new_eval_reset = True
            else:
                reset_eval_form()
            st.rerun()

        if st.session_state.get("confirm_new_eval_reset", False):
            st.warning("작성 중인 검토 문장이 있습니다. 새 검토를 시작하면 현재 입력 내용이 초기화됩니다.")
            col_e1, col_e2 = st.columns(2)
            with col_e1:
                if st.button("새 검토 계속", type="primary", key="btn_confirm_eval_reset"):
                    reset_eval_form()
                    st.rerun()
            with col_e2:
                if st.button("취소", key="btn_cancel_eval_reset"):
                    st.session_state.confirm_new_eval_reset = False
                    st.rerun()

        st.subheader("검토 목록")
        
        # 1. 저장된 기존 검토 목록 (순서 및 위치 고정)
        for idx, rec in enumerate(st.session_state.eval_records):
            btn_label = f"{rec.get('title', f'검토 {idx+1}')}"
            is_active = (not st.session_state.is_new_eval) and (idx == st.session_state.current_eval_idx)
            if st.button(
                btn_label, 
                key=f"eval_btn_{idx}", 
                use_container_width=True, 
                type="primary" if is_active else "secondary"
            ):
                st.session_state.is_new_eval = False
                st.session_state.current_eval_idx = idx
                st.session_state.draft_eval_text = rec.get("target_text", "")
                st.session_state.eval_target_text = rec.get("target_text", "")
                st.session_state.eval_text_widget = rec.get("target_text", "")
                st.session_state.eval_result = rec.get("result", "")
                st.session_state.confirm_new_eval_reset = False
                if rec.get("model") and rec.get("model") in MODEL_MAP:
                    st.session_state.selected_model_name = rec.get("model")
                eval_key = f"uploader_eval_{st.session_state.uploader_key_eval}"
                if eval_key in st.session_state:
                    del st.session_state[eval_key]
                st.session_state.uploader_key_eval += 1
                st.session_state.attachments["eval"] = []
                st.rerun()

        # 2. 임시 신규 검토 항목 (목록 맨 아래에 파란색 활성화 표시)
        if st.session_state.is_new_eval:
            st.button("새 검토", key="temp_new_eval_btn", use_container_width=True, type="primary")

    st.markdown("---")
    
    # 일괄 엑셀 다운로드 (생기부 작성 모드에서만 노출)
    if mode == "생기부 작성":
        if st.session_state.student_records:
            data_list = []
            for r in st.session_state.student_records:
                c_cnt, b_cnt = calculate_neis_bytes(r["draft"])
                data_list.append({
                    "학번": r["id_val"],
                    "작성영역": r["record_type"],
                    "목표바이트": r.get("target_bytes", 1500),
                    "관찰내용 및 키워드": r["memo"],
                    "생기부내용 (초안)": r["draft"],
                    "글자수": c_cnt,
                    "NEIS 바이트": b_cnt
                })
            df_bulk = pd.DataFrame(data_list)
            excel_buffer = io.BytesIO()
            with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
                df_bulk.to_excel(writer, index=False, sheet_name='생기부_일괄내역')
            excel_data = excel_buffer.getvalue()
            
            st.download_button(
                label="전체 학생 기록 다운로드 (Excel)",
                data=excel_data,
                file_name="전체생기부기록.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True
            )
        else:
            st.button("일괄 엑셀 다운로드 (작성 데이터 없음)", disabled=True, use_container_width=True)
        st.markdown("---")

    st.caption(f"누적 사용 토큰: **{st.session_state.total_tokens_used:,} Tokens**")
    
    if st.button("로그아웃", use_container_width=True):
        st.session_state.authenticated = False
        st.session_state.access_mode = None
        st.rerun()

# 6. 메인 동작 영역
if mode == "일반 챗봇":
    curr_idx = st.session_state.current_chat_idx if st.session_state.current_chat_idx is not None else 0
    current_chat = st.session_state.chat_sessions[curr_idx]
    
    # 1. Streamlit 네이티브 스크롤 컨테이너 (고정 높이 및 자체 스크롤 제공)
    with st.container(key="chat_history_area", border=False):
        # 대화 히스토리 화면 출력
        for msg in current_chat["messages"]:
            with st.chat_message(msg["role"]):
                # 사용자 첨부 원문은 API 처리용으로 보존하되 화면에는 질문/첨부명만 간결하게 표시
                if msg.get("display_text") is not None:
                    st.write(msg.get("display_text") or "첨부자료 분석 요청")
                elif isinstance(msg["content"], list):
                    for part in msg["content"]:
                        if part["type"] == "text":
                            st.write(part["text"])
                        elif part["type"] == "image_url":
                            st.image(part["image_url"]["url"], use_container_width=False)
                else:
                    st.write(msg["content"])
                if msg.get("excel_file"):
                    ef = msg.get("excel_file") or {}
                    if ef.get("data") and ef.get("filename"):
                        st.download_button("Excel 다운로드", data=ef["data"], file_name=ef["filename"], mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", key=f"excel_hist_{curr_idx}_{id(msg)}")
                if "tokens" in msg:
                    in_tok = msg.get("input_tokens", 0)
                    out_tok = msg.get("output_tokens", 0)
                    if in_tok or out_tok:
                        st.caption(f"이번 요청: 입력 {in_tok:,} · 출력 {out_tok:,} · 합계 {msg['tokens']:,} Tokens")
                    else:
                        st.caption(f"소모 토큰: {msg['tokens']:,} Tokens")

        # 새 질문 전송 후 답변 생성 상태인 경우 대화창 내부에서 assistant 플레이스홀더 생성 및 API 호출
        if st.session_state.get("should_generate_chat", False):
            st.session_state.should_generate_chat = False
            with st.chat_message("assistant"):
                answer_placeholder = st.empty()
                with answer_placeholder.container():
                    st.caption(f"[{selected_model_name}] 답변 생성 중...")
                
                if not check_openrouter_model_availability(selected_model):
                    answer_placeholder.error("현재 선택한 AI 모델을 OpenRouter에서 사용할 수 없습니다. 관리자에게 모델 설정을 확인해 주세요.")
                else:
                    client = get_openrouter_client()
                    if not client:
                        answer_placeholder.error(get_api_unavailable_message())
                    else:
                        try:
                            current_user_message = current_chat["messages"][-1]
                            api_messages = [{"role": "system", "content": build_chat_system_prompt(current_chat, current_user_message)}]

                            # 최근 대화만 문맥으로 유지하고, 과거 첨부 원문/Base64는 반복 전송하지 않음
                            recent_history = current_chat["messages"][:-1][-CHAT_HISTORY_MAX_MESSAGES:]
                            for m in recent_history:
                                api_messages.append({
                                    "role": m["role"],
                                    "content": compact_history_content(m) if m["role"] == "user" else m["content"]
                                })

                            # 현재 질문의 첨부자료만 원본 형태로 전송
                            api_messages.append({"role": "user", "content": current_chat["messages"][-1]["content"]})
                            
                            response = client.chat.completions.create(
                                model=selected_model,
                                messages=api_messages
                            )
                            raw_res = response.choices[0].message.content or ""
                            visible_res, excel_spec, file_cache = parse_hidden_chat_payloads(raw_res)
                            input_tokens, output_tokens, tokens_count = extract_usage_tokens(response)
                            if tokens_count > 0:
                                st.session_state.total_tokens_used += tokens_count
                            
                            if file_cache:
                                current_chat["source_cache_text"] = file_cache[:CHAT_FILE_CACHE_MAX_CHARS]
                            excel_bytes = None
                            excel_filename = None
                            if excel_spec:
                                excel_bytes, excel_filename = build_excel_bytes(excel_spec)
                                spec_cache = excel_spec_to_cache(excel_spec)
                                if spec_cache:
                                    current_chat["source_cache_text"] = spec_cache
                            display_res = (visible_res or raw_res)
                            if excel_bytes and not visible_res:
                                display_res = f"첨부자료를 정리해 **{excel_filename}** 파일을 만들었습니다. 아래에서 다운로드할 수 있습니다."
                            answer_placeholder.markdown(display_res)
                            if excel_bytes:
                                st.download_button("Excel 다운로드", data=excel_bytes, file_name=excel_filename, mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", key=f"excel_now_{curr_idx}_{len(current_chat['messages'])}")
                            if tokens_count > 0:
                                if input_tokens or output_tokens:
                                    st.caption(f"이번 요청: 입력 {input_tokens:,} · 출력 {output_tokens:,} · 합계 {tokens_count:,} Tokens")
                                else:
                                    st.caption(f"소모 토큰: {tokens_count:,} Tokens")
                            new_msg = {"role": "assistant", "content": display_res}
                            if excel_bytes:
                                new_msg["excel_file"] = {"data": excel_bytes, "filename": excel_filename}
                            if tokens_count > 0:
                                new_msg["tokens"] = tokens_count
                                new_msg["input_tokens"] = input_tokens
                                new_msg["output_tokens"] = output_tokens
                            current_chat["messages"].append(new_msg)
                            st.session_state.force_chat_scroll = True

                            # 사이드바의 누적 토큰을 같은 요청 직후 즉시 갱신
                            st.rerun()
                        except Exception as e:
                            answer_placeholder.error(f"AI 응답 생성 실패: {str(e)}")

        # 긴 답변의 마지막 토큰 사용량이 하단 고정 입력창에 가려지지 않도록
        # 스크롤 영역 끝에 항상 안전 여백을 둔다.
        st.markdown('<div class="chat-bottom-safe-space"></div>', unsafe_allow_html=True)

        import streamlit.components.v1 as components
        force_scroll_json = json.dumps(bool(st.session_state.get("force_chat_scroll", False)))
        components.html(f"""
            <script>
            function autoScrollChat() {{
                try {{
                    const force = {force_scroll_json};
                    const targets = window.parent.document.querySelectorAll('.st-key-chat_history_area, div[key="chat_history_area"]');
                    targets.forEach(function(container) {{
                        if (!container) return;
                        const scrollEl = container.querySelector('[data-testid="stScrollArea"]') || container;
                        const distFromBottom = scrollEl.scrollHeight - scrollEl.scrollTop - scrollEl.clientHeight;
                        if (force || distFromBottom < 300) {{
                            scrollEl.scrollTop = scrollEl.scrollHeight + 300;
                        }}
                    }});
                }} catch(e) {{}}
            }}
            autoScrollChat();
            setTimeout(autoScrollChat, 50);
            setTimeout(autoScrollChat, 180);
            </script>
        """, height=0)
        st.session_state.force_chat_scroll = False

    # 2. Streamlit 1.61.1 공식 하단 고정 컨테이너 st.bottom 활용
    with st.bottom:
        with st.container(key="chat_composer_area"):
            # paste listener callback
            paste_key = "paste_listener_chat"
            def handle_pasted_image():
                state = st.session_state.get(paste_key)
                if state and hasattr(state, "pasted_image") and state.pasted_image:
                    img_data = state.pasted_image
                    add_attachment(
                        name=img_data["name"],
                        mime_type=img_data["type"],
                        data_base64=img_data["data"],
                        size_bytes=img_data["size"],
                        source="paste",
                        file_type="image",
                        scope="chat"
                    )
            paste_listener(key=paste_key, on_pasted_image=handle_pasted_image)

            # ChatGPT 스타일 콤팩트 대기 첨부 칩스 (Composer 내부에만 렌더링)
            current_attachments = st.session_state.attachments["chat"]
            visible_chat_attachments = [(idx, item) for idx, item in enumerate(current_attachments) if not item.get("hidden", False)]
            if visible_chat_attachments:
                with st.container(key="chat_pending_chips"):
                    max_cols = min(len(visible_chat_attachments), 6)
                    cols = st.columns(max_cols)
                    for visible_idx, (original_idx, item) in enumerate(visible_chat_attachments):
                        with cols[visible_idx % max_cols]:
                            if item["type"] == "image":
                                st.image(item["data"], width=80)
                            else:
                                st.caption(f"문서: {item['name'][:10]}")
                            if st.button("×", key=f"del_chat_att_{original_idx}", help=f"{item['name']} 삭제"):
                                delete_attachment_group("chat", item.get("hash"))
                                st.rerun()

            # native st.chat_input (단 1개만 렌더링)
            user_input = st.chat_input(
                "Chat PSDongSung에게 물어보기 (파일 첨부 및 Ctrl+V 캡처 지원)",
                key="chat_input_widget",
                accept_file="multiple",
                file_type=["png","jpg","jpeg","webp","pdf","hwp","hwpx","xlsx","xls","docx","pptx"]
            )

    if user_input:
        prompt = ""
        if hasattr(user_input, "text"):
            prompt = user_input.text or ""
        elif isinstance(user_input, dict):
            prompt = user_input.get("text", "")
        elif isinstance(user_input, str):
            prompt = user_input

        submitted_files = []
        if hasattr(user_input, "files") and user_input.files:
            submitted_files = user_input.files
        elif isinstance(user_input, dict) and user_input.get("files"):
            submitted_files = user_input.get("files")

        file_meta = []
        if submitted_files:
            for f in submitted_files:
                file_bytes = f.getvalue()
                file_hash = calculate_bytes_hash(file_bytes)
                lower_name = f.name.lower()
                if lower_name.endswith(('.png', '.jpg', '.jpeg', '.webp')):
                    import base64
                    b64_data = f"data:{f.type};base64," + base64.b64encode(file_bytes).decode("utf-8")
                    add_attachment(f.name, f.type, b64_data, f.size, "upload", "image", scope="chat")
                    file_meta.append({"name": f.name, "kind": "image", "hash": file_hash})
                elif lower_name.endswith('.pdf'):
                    total_pages = get_pdf_page_count(file_bytes)
                    vision_pages = add_pdf_attachment_with_vision(f, "chat", file_bytes, file_hash)
                    sent_pages = total_pages if total_pages else vision_pages
                    file_meta.append({"name": f.name, "kind": "pdf", "hash": file_hash, "total_pages": total_pages, "sent_pages": sent_pages})
                else:
                    text_content = parse_uploaded_file(f)
                    add_document_attachment(f.name, f.type, text_content, f.size, file_hash, scope="chat")
                    file_meta.append({"name": f.name, "kind": "document", "hash": file_hash, "description": "문서 본문을 추출하여 전달"})
        attachment_names = [item.get("name", "첨부파일") for item in st.session_state.attachments["chat"] if not item.get("hidden", False)]
        metadata_context = build_file_metadata_context(file_meta)
        prompt_for_api = prompt
        if metadata_context:
            prompt_for_api = (prompt_for_api + "\n\n" if prompt_for_api else "") + metadata_context
        user_content = compile_api_payload(prompt_for_api, selected_model_name, scope="chat")

        attachment_note = ""
        if attachment_names:
            attachment_note = "첨부: " + ", ".join(attachment_names)
        history_text = prompt.strip()
        if attachment_note:
            history_text += ("\n" if history_text else "") + f"[{attachment_note}]"
        if not history_text:
            history_text = "첨부자료 분석 요청"

        current_chat["messages"].append({
            "role": "user",
            "content": user_content,
            "history_text": history_text,
            "display_text": history_text,
            "file_meta": file_meta,
            "excel_intent": wants_excel_output(prompt)
        })
        
        if current_chat["title"] == "새 대화":
            current_chat["title"] = prompt[:12] + "..." if len(prompt) > 12 else (prompt if prompt else "첨부자료 분석")
            
        st.session_state.attachments["chat"] = []
        chat_key = f"uploader_chat_{st.session_state.uploader_key_chat}"
        if chat_key in st.session_state:
            del st.session_state[chat_key]
        st.session_state.uploader_key_chat += 1

        st.session_state.should_generate_chat = True
        st.session_state.force_chat_scroll = True
        st.rerun()

elif mode == "생기부 작성":
    # 2. Streamlit 네이티브 스크롤 컨테이너 (CSS dvh 동적 가용 높이 제어)
    with st.container(border=False, key="content_area"):
        curr_std_idx = st.session_state.current_student_idx
        if st.session_state.is_new_student or curr_std_idx is None:
            st.caption("**현재 상태: 새 학생 작성 중**")
        else:
            curr_std = st.session_state.student_records[curr_std_idx]
            name_part = f" {curr_std.get('name', '')}" if curr_std.get('name') else ""
            st.caption(f"**현재 작성: {curr_std['id_val']}{name_part} · {curr_std.get('record_type', '교과세특')}**")

        with st.expander("공통 참조자료", expanded=False):
            st.markdown(
                """
                **여러 학생에게 공통으로 적용할 자료를 넣어두면 좋습니다.**  
                - 교육과정 **성취기준** 및 내용 요소
                - 성취수준, **평가기준** 및 채점 기준표
                - 수행평가 계획서, 탐구 과제 안내, 수업 활동 자료
                - 과목에서 강조할 핵심역량 및 교과별 관찰 관점
                - 선생님의 **개별 작성 프롬프트**, 문체·어조·표현 규칙
                - 학교 또는 부서의 생기부 기재 원칙, 금지 표현, 유의사항
                - 참고할 만한 우수 작성 예시나 기존 기록 양식

                학생 개인의 보고서·사진·관찰자료는 아래 **현재 학생 첨부자료**에 넣는 것을 권장합니다.
                """
            )
            global_ref_files = st.file_uploader(
                "공통 참조 파일 업로드 (PDF, HWP 등)",
                accept_multiple_files=True,
                key="global_ref_uploader",
                label_visibility="collapsed"
            )
            global_ref_text = ""
            global_ref_hashes = []
            if global_ref_files:
                ref_texts = []
                for g_file in global_ref_files:
                    g_bytes = g_file.getvalue()
                    g_hash = calculate_bytes_hash(g_bytes)
                    global_ref_hashes.append(g_hash)
                    ref_text = parse_uploaded_file(g_file)
                    ref_texts.append(f"\n\n[공통참조문서: {g_file.name}]\n" + ref_text)
                    if g_file.name.lower().endswith('.pdf') and ref_text.startswith(PDF_IMAGE_MARKER):
                        # 렌더링 결과는 st.cache_data에 저장되어 같은 PDF를 화면 rerun 때 다시 변환하지 않습니다.
                        pages, _ = render_pdf_pages_as_images(g_bytes)
                        for page_no, data_url, image_size in pages:
                            add_attachment(
                                name=f"공통참조 {g_file.name} · {page_no}쪽",
                                mime_type="image/jpeg",
                                data_base64=data_url,
                                size_bytes=image_size,
                                source="global_pdf_page",
                                file_type="image",
                                scope="student",
                                hidden=True,
                                parent_hash=f"global_{g_hash}"
                            )
                global_ref_text = "".join(ref_texts)
                st.success(f"총 {len(global_ref_files)}개의 공통 참조 문서가 반영되었습니다. 원본 분석은 최초 1회만 수행합니다.")
            else:
                global_ref_hashes = []

        if "student_id_widget" not in st.session_state:
            st.session_state["student_id_widget"] = st.session_state.get("draft_student_id", "")
        if "student_record_type_widget" not in st.session_state:
            st.session_state["student_record_type_widget"] = st.session_state.get("draft_student_record_type", "교과세특")
        if "student_target_bytes_widget" not in st.session_state:
            st.session_state["student_target_bytes_widget"] = int(st.session_state.get("draft_student_target_bytes", 1500))
        if "student_memo_widget" not in st.session_state:
            st.session_state["student_memo_widget"] = st.session_state.get("draft_student_memo", "")

        # 새 학생 작성 화면은 작성 영역의 공식 기본 Byte에서 시작합니다.
        if st.session_state.is_new_student and st.session_state.current_student_idx is None:
            _rt_for_default = st.session_state.get("draft_student_record_type", "교과세특") or "교과세특"
            if st.session_state.get("_new_student_default_bytes_applied") != _rt_for_default:
                _default_b = get_default_target_bytes(_rt_for_default)
                st.session_state.student_target_bytes_widget = _default_b
                st.session_state.draft_student_target_bytes = _default_b
                st.session_state._new_student_default_bytes_applied = _rt_for_default

        record_options = ["교과세특", "동아리활동", "자율활동", "진로활동", "행동특성 및 종합의견"]
        current_record_default = st.session_state.get("draft_student_record_type", "교과세특")
        if current_record_default not in record_options:
            current_record_default = "교과세특"

        def _sync_target_bytes_to_record_type():
            selected = st.session_state.get("student_record_type_segmented", "교과세특") or "교과세특"
            default_bytes = get_default_target_bytes(selected)
            st.session_state.student_target_bytes_widget = default_bytes
            st.session_state.draft_student_target_bytes = default_bytes
            st.session_state._new_student_default_bytes_applied = selected

        record_type = st.segmented_control(
            "작성 영역",
            record_options,
            default=current_record_default,
            key="student_record_type_segmented",
            selection_mode="single",
            on_change=_sync_target_bytes_to_record_type
        ) or "교과세특"
        st.session_state["draft_student_record_type"] = record_type

        col_std1, col_std2 = st.columns([1, 1])
        with col_std1:
            student_id = st.text_input("학번", key="student_id_widget", placeholder="예: 10101")
            st.session_state["draft_student_id"] = student_id
        with col_std2:
            target_bytes = int(st.number_input(
                "목표 바이트",
                min_value=300,
                max_value=5000,
                step=50,
                key="student_target_bytes_widget",
                help=f"{record_type} 기본값은 {get_default_target_bytes(record_type):,} Byte입니다. 필요하면 직접 조정할 수 있습니다."
            ))
            st.session_state["draft_student_target_bytes"] = target_bytes
            
        student_memo = st.text_area("학생 관찰 내용 및 키워드", key="student_memo_widget", height=95, placeholder="수업 참여도, 수행평가 과정, 특기사항 등 입력")
        st.session_state["draft_student_memo"] = student_memo
        
        render_attachments_panel(uploader_key=f"uploader_std_{st.session_state.uploader_key_std}", scope="student")

        student_attachment_text = get_scope_document_text("student")
        is_structured_table_source = "[STRUCTURED_TABLE_TSV]" in student_attachment_text

        # v15: XLSX/CSV 헤더(예: 이름, 희망진로)를 학생 이름으로 오인식하지 않도록
        # 구조화 표는 추측형 정규식 식별을 사용하지 않고 학번/이름 열만 신뢰합니다.
        if is_structured_table_source:
            inferred_id, inferred_name = "", ""
            batch_candidates = []
        else:
            inferred_id, inferred_name = infer_student_identity(student_attachment_text)
            inferred_name = clean_student_name(inferred_name, inferred_id)
            batch_candidates = extract_student_candidates(student_attachment_text)

        effective_student_id = student_id.strip() or inferred_id
        effective_student_name = inferred_name

        if len(batch_candidates) >= 2:
            st.info(f"첨부자료에서 **학생 후보 {len(batch_candidates)}명**을 감지했습니다. 1회 분석 후 학생별로 정리합니다.")
        elif (not is_structured_table_source) and (inferred_id or inferred_name):
            identity_text = " ".join(x for x in [inferred_id, inferred_name] if x)
            st.info(f"첨부자료에서 인식한 학생: **{identity_text}**")

        # --------------------------------------------------
        # v13: 원본 첨부자료는 최초 1회만 분석하고 이후에는 요약 캐시만 사용
        # --------------------------------------------------
        candidate_hint = ", ".join(f"{sid} {name}" for sid, name in batch_candidates[:50])
        source_signature = build_student_source_signature(global_ref_hashes)
        has_student_sources = bool(source_signature)
        source_analysis = st.session_state.student_source_analysis_cache.get(source_signature) if source_signature else None

        # XLSX/CSV처럼 이미 행과 열이 명확한 자료는 AI에게 1만 토큰 이상 보내지 않고 로컬에서 학생별로 직접 분리
        if has_student_sources and not source_analysis:
            structured_analysis = build_structured_student_analysis_from_text(student_attachment_text)
            if structured_analysis and structured_analysis.get("students"):
                st.session_state.student_source_analysis_cache[source_signature] = structured_analysis
                source_analysis = structured_analysis
                st.session_state["_last_source_analysis_tokens"] = {"input": 0, "output": 0, "total": 0}
                st.session_state["_source_analysis_mode"] = "structured_local"

        if has_student_sources:
            if source_analysis:
                analyzed_students = source_analysis.get("students", []) or []
                if st.session_state.get("_source_analysis_mode") == "structured_local":
                    st.success(f"구조화된 표를 직접 분석했습니다 · 학생 {len(analyzed_students)}명 감지 · 자료분석 AI 토큰 0")
                else:
                    st.success(f"첨부자료 분석 캐시 사용 중 · 학생 {len(analyzed_students)}명 감지 · 원본 PDF/이미지는 다시 보내지 않습니다.")
            else:
                st.info("첨부자료는 최초 1회만 AI가 분석합니다. 이후 초안 생성·바이트 조정·추가 수정에는 원본 PDF/이미지를 다시 보내지 않습니다.")
                if st.button("첨부자료 1회 분석", key="btn_analyze_student_sources", use_container_width=True, disabled=st.session_state.get("is_analyzing_student_sources", False)):
                    st.session_state.is_analyzing_student_sources = True
                    st.session_state._run_source_analysis_once = True
                    st.session_state["_pending_action_after_source_analysis"] = ""
                    st.rerun()

        if (st.session_state.get("is_analyzing_student_sources", False)
                and st.session_state.get("_run_source_analysis_once", False)):
            st.session_state._run_source_analysis_once = False
            client = get_openrouter_client()
            if not client:
                st.session_state.is_analyzing_student_sources = False
                st.error(get_api_unavailable_message())
            else:
                with st.spinner(f"[{selected_model_name}] 첨부자료를 최초 1회 분석 중... 이후에는 이 요약을 재사용합니다."):
                    try:
                        analysis, ain, aout, atotal = analyze_student_sources_once(
                            client, selected_model, selected_model_name,
                            global_ref_text=global_ref_text,
                            candidate_hint=candidate_hint,
                        )
                        st.session_state.student_source_analysis_cache[source_signature] = analysis
                        st.session_state["_source_analysis_mode"] = "ai"
                        if atotal > 0:
                            st.session_state.total_tokens_used += atotal
                        st.session_state["_last_source_analysis_tokens"] = {"input": ain, "output": aout, "total": atotal}
                        st.session_state.is_analyzing_student_sources = False
                        pending_action = st.session_state.pop("_pending_action_after_source_analysis", "")
                        if pending_action == "single":
                            st.session_state.is_generating_draft = True
                            st.session_state._run_draft_once = True
                    except Exception as e:
                        st.session_state.is_analyzing_student_sources = False
                        st.error(f"첨부자료 분석 중 오류 발생: {str(e)}")
                    else:
                        st.rerun()

        # 분석 완료 후 핵심 추출 내용과 여러 학생 선택 UI 제공
        source_analysis = st.session_state.student_source_analysis_cache.get(source_signature) if source_signature else None
        analyzed_students = (source_analysis or {}).get("students", []) or []
        if source_analysis and analyzed_students and not effective_student_id:
            first_detected = analyzed_students[0]
            effective_student_id = str(first_detected.get("id", "")).strip()
            effective_student_name = str(first_detected.get("name", "")).strip()
            if len(analyzed_students) == 1 and (effective_student_id or effective_student_name):
                identity_text = " ".join(x for x in [effective_student_id, effective_student_name] if x)
                st.info(f"1회 분석에서 인식한 학생: **{identity_text}**")
            elif len(analyzed_students) >= 2:
                st.success(f"1회 분석 완료 · **학생 {len(analyzed_students)}명**을 학생별로 분리했습니다.")

        if source_analysis:
            with st.expander("첨부자료에서 추출한 핵심 내용", expanded=False):
                if analyzed_students:
                    for item in analyzed_students[:30]:
                        who = " ".join(x for x in [item.get("id", ""), item.get("name", "")] if x).strip() or "학생 미식별"
                        st.markdown(f"**{who}** — {item.get('summary', '')}")
                elif source_analysis.get("source_summary"):
                    st.write(source_analysis.get("source_summary"))
                last_t = st.session_state.get("_last_source_analysis_tokens", {})
                if last_t.get("total"):
                    st.caption(f"최초 자료 분석: 입력 {last_t.get('input',0):,} · 출력 {last_t.get('output',0):,} · 합계 {last_t.get('total',0):,} Tokens")

        if len(analyzed_students) >= 2:
            labels = []
            label_to_student = {}
            for idx, item in enumerate(analyzed_students):
                sid = str(item.get("id", "")).strip()
                name = str(item.get("name", "")).strip()
                label = " ".join(x for x in [sid, name] if x).strip() or f"학생 {idx+1}"
                # 중복 라벨 방지
                if label in label_to_student:
                    label = f"{label} ({idx+1})"
                labels.append(label)
                label_to_student[label] = item
            with st.expander(f"여러 학생 자료 감지 · {len(analyzed_students)}명", expanded=True):
                st.caption("대상을 선택해 한 번에 생성할 수 있습니다. 일괄 생성은 토큰과 시간을 아끼기 위해 원본 파일을 보내지 않고 학생별 요약만 사용하며, 자동 Byte 재보정은 생략합니다.")
                select_key = f"batch_select_{source_signature[:10]}"
                selected_labels = st.multiselect("초안 생성 대상", labels, default=labels, key=select_key)
                if st.button("선택 학생 초안 일괄 생성", key="btn_batch_students_v11", type="primary", use_container_width=True, disabled=not selected_labels):
                    st.session_state["_pending_batch_students"] = [label_to_student[x] for x in selected_labels]
                    st.session_state.is_generating_batch_students = True
                    st.session_state._run_batch_once = True
                    st.rerun()

                failed = st.session_state.get("batch_failed_students", [])
                if failed:
                    st.warning(f"이전 일괄 생성에서 {len(failed)}명 실패: " + ", ".join(" ".join([str(x.get('id','')), str(x.get('name',''))]).strip() for x in failed))
                    if st.button("실패 학생만 다시 생성", key="btn_retry_failed_batch", use_container_width=True):
                        st.session_state["_pending_batch_students"] = failed
                        st.session_state.is_generating_batch_students = True
                        st.session_state._run_batch_once = True
                        st.rerun()

        if st.session_state.get("is_generating_batch_students", False) and st.session_state.get("_run_batch_once", False):
            st.session_state._run_batch_once = False
            pending_students = st.session_state.get("_pending_batch_students", []) or []
            client = get_openrouter_client()
            if not client:
                st.session_state.is_generating_batch_students = False
                st.error(get_api_unavailable_message())
            elif pending_students:
                progress = st.progress(0.0, text=f"학생 초안 준비 중 · 0/{len(pending_students)}")
                status_box = st.empty()
                successes = []
                failures = []
                # 4명씩 텍스트 요약만 묶어 호출: 이미지 재전송 없이 속도/비용 절약
                chunk_size = 4
                common_summary = (source_analysis or {}).get("common_summary", "")
                try:
                    for chunk_start in range(0, len(pending_students), chunk_size):
                        chunk = pending_students[chunk_start:chunk_start + chunk_size]
                        compact_items = []
                        for item in chunk:
                            compact_items.append({
                                "id": str(item.get("id", "")).strip(),
                                "name": str(item.get("name", "")).strip(),
                                "summary": str(item.get("summary", "")).strip(),
                            })
                        batch_text = (
                            f"작성 영역: {record_type}\n목표: 학생별 {target_bytes} Byte 이하, 권장 90~100%\n"
                            f"공통 기준 요약: {common_summary}\n"
                            f"교사 공통 메모: {student_memo or '없음'}\n\n"
                            "다음 학생별 요약을 서로 섞지 말고 생기부 초안으로 작성하세요. 자료에 없는 사실은 만들지 마세요. "
                            "학생 이름을 주어로 시작하지 말고 문장 끝은 '~함', '~임' 중심으로 작성하세요. "
                            "반드시 JSON 배열만 출력하세요. 각 항목은 id, name, draft 필드입니다.\n"
                            + json.dumps(compact_items, ensure_ascii=False)
                        )
                        resp = client.chat.completions.create(
                            model=selected_model,
                            messages=[
                                {"role": "system", "content": "학교생활기록부 작성 전문가입니다. 학생별 근거를 섞지 말고 JSON 외 텍스트를 출력하지 마세요."},
                                {"role": "user", "content": batch_text},
                            ],
                        )
                        rin, rout, rtotal = extract_usage_tokens(resp)
                        if rtotal > 0:
                            st.session_state.total_tokens_used += rtotal
                        raw = re.sub(r"^```(?:json)?\s*|\s*```$", "", resp.choices[0].message.content.strip(), flags=re.I | re.S).strip()
                        try:
                            results = json.loads(raw)
                            if not isinstance(results, list):
                                raise ValueError("배열 아님")
                        except Exception:
                            results = []
                        by_id = {str(x.get("id", "")).strip(): x for x in results if isinstance(x, dict)}
                        for item in chunk:
                            sid = str(item.get("id", "")).strip()
                            name = str(item.get("name", "")).strip()
                            result_item = by_id.get(sid)
                            if result_item is None and name:
                                result_item = next((x for x in results if isinstance(x, dict) and str(x.get("name", "")).strip() == name), None)
                            draft = str((result_item or {}).get("draft", "")).strip()
                            if not draft:
                                failures.append(item)
                            else:
                                share_count = max(1, len(chunk))
                                rec = {
                                    "info": f"학번: {sid}" + (f", 이름: {name}" if name else ""),
                                    "id_val": sid or name,
                                    "name": name,
                                    "memo": "첨부자료 1회 분석 요약 기반 일괄 생성",
                                    "draft": draft,
                                    "record_type": record_type,
                                    "target_bytes": target_bytes,
                                    # 일괄 호출 토큰은 학생 수로 나눈 참고값을 기록하고, 실제 전체 사용량은 사이드바 누적값에 정확히 반영합니다.
                                    "tokens": int(round(rtotal / share_count)),
                                    "input_tokens": int(round(rin / share_count)),
                                    "output_tokens": int(round(rout / share_count)),
                                    "source_signature": source_signature,
                                }
                                # 같은 학번이 있으면 덮어쓰고, 없으면 추가
                                existing_idx = next((i for i, r in enumerate(st.session_state.student_records) if sid and str(r.get("id_val", "")) == sid), None)
                                if existing_idx is not None:
                                    st.session_state.student_records[existing_idx] = rec
                                else:
                                    st.session_state.student_records.append(rec)
                                successes.append(item)
                            done = len(successes) + len(failures)
                            progress.progress(min(1.0, done / max(1, len(pending_students))), text=f"학생 초안 생성 · {done}/{len(pending_students)}")
                            status_box.caption(f"최근 처리: {sid} {name}".strip())
                    st.session_state.batch_failed_students = failures
                    if successes:
                        first_sid = str(successes[0].get("id", "")).strip()
                        idx = next((i for i, r in enumerate(st.session_state.student_records) if first_sid and str(r.get("id_val", "")) == first_sid), len(st.session_state.student_records)-1)
                        st.session_state.current_student_idx = max(0, idx)
                        st.session_state.is_new_student = False
                    st.session_state["_pending_batch_students"] = []
                    st.session_state.is_generating_batch_students = False
                    if failures:
                        st.warning(f"{len(successes)}명 완료, {len(failures)}명 실패했습니다. 실패 학생만 다시 생성할 수 있습니다.")
                    else:
                        st.success(f"{len(successes)}명의 학생 초안을 생성했습니다.")
                except Exception as e:
                    st.session_state.is_generating_batch_students = False
                    st.error(f"학생 일괄 생성 중 오류 발생: {str(e)}")
                else:
                    st.rerun()

        is_draft_busy = st.session_state.get("is_generating_draft", False)
        if st.button("초안 생성 중..." if is_draft_busy else "초안 생성", type="primary", use_container_width=True, disabled=is_draft_busy):
            if effective_student_id:
                if has_student_sources and not source_analysis:
                    # 원본이 아직 분석되지 않았다면 최초 1회 분석을 먼저 수행하고 자동으로 초안 생성으로 이어갑니다.
                    st.session_state.is_analyzing_student_sources = True
                    st.session_state._run_source_analysis_once = True
                    st.session_state["_pending_action_after_source_analysis"] = "single"
                else:
                    st.session_state.is_generating_draft = True
                    st.session_state._run_draft_once = True
                st.rerun()
            else:
                st.warning("학번을 입력하거나, 학번이 포함된 학생 자료를 첨부해 주세요.")

        if st.session_state.get("is_generating_draft", False) and st.session_state.get("_run_draft_once", False):
            # 이 실행에서 요청을 즉시 소비합니다. Stop 후 rerun되어도 자동 재시작하지 않습니다.
            st.session_state._run_draft_once = False
            resolved_student_id = effective_student_id
            resolved_student_name = effective_student_name
            info_str = f"학번: {resolved_student_id}" + (f", 이름: {resolved_student_name}" if resolved_student_name else "")
            if not check_openrouter_model_availability(selected_model):
                st.error("현재 선택한 AI 모델을 OpenRouter에서 사용할 수 없습니다. 관리자에게 모델 설정을 확인해 주세요.")
                st.session_state.is_generating_draft = False
            else:
                client = get_openrouter_client()
                if not client:
                    st.error(get_api_unavailable_message())
                    st.session_state.is_generating_draft = False
                else:
                    with st.spinner(f"[{selected_model_name}] {resolved_student_id} {resolved_student_name} 초안 생성 중..."):
                        try:
                            system_prompt = (
                                "당신은 대한민국 교육부 및 시도교육청의 학교생활기록부 작성 전문가입니다.\n"
                                "2022 개정 교육과정 기준 및 학교생활기록부 기재요령에 맞추어 다음 작성 지침을 '엄격히' 준수하세요:\n\n"
                                "1. [문장 시작 규칙 - 필수]: 문장 시작 시 '000 학생은', 'OOO 학생은', 'OO은' 같은 이름이나 학생 주어를 절대로 사용하지 마세요. 곧바로 학생의 학업적 성취 특성이나 활동 내용으로 문장을 바로 시작하세요.\n\n"
                                "2. [표현 제한 - 필수]: 영어 단어, 알파벳, 괄호 안 영문 병기(예: English) 및 중간점(·) 등의 특수문자를 절대로 사용하지 마세요. 모든 개념과 단어는 정돈된 표준 한글로만 작성하세요.\n\n"
                                "3. [공통 참조 자료 반영]: 제시된 [공통 교육과정/성취기준 참조 자료]의 성취기준 및 교과역량을 적극 반영하여 작성하세요.\n\n"
                                "4. [영역별 작성 지침]:\n"
                                f"   - 현재 작성 영역: {record_type}\n"
                                "   - 교과세특인 경우: 성취수준 + 수행 과정 및 결과(과제, 분석내용, 도구) + 교과 핵심역량 + 교사 총평 구조 포함.\n"
                                "   - 창체/행발인 경우: 행동 특성, 공동체 의식, 주도적 활동 및 변화 모습을 구체적으로 진술.\n\n"
                                "5. [어조 및 분량]:\n"
                                "   - 문장 끝은 반드시 '~함', '~임' 어조로만 작성하세요.\n"
                                f"   - 목표 분량은 {target_bytes} Byte입니다. 가급적 목표의 90~100% 범위로 작성하고 {target_bytes} Byte를 넘기지 마세요.\n"
                                "   - 분량을 맞추기 위해 의미 없는 반복 표현을 넣지 말고, 구체적인 수행 과정과 교사의 관찰 근거를 우선 보강하세요.\n\n"
                                "6. [기재 금지어]: 대회, 수상, 외부 기관명, 공인어학성적, 사교육 관련 내용 절대 언급 금지."
                            )
                            
                            # v13: 최초 분석 이후에는 원본 PDF/HWPX/이미지를 절대 재전송하지 않고 요약만 사용합니다.
                            active_analysis = st.session_state.student_source_analysis_cache.get(source_signature) if source_signature else None
                            source_summary_text = find_student_source_summary(active_analysis, resolved_student_id, resolved_student_name)
                            prompt_base = f"[작성 정보]: {info_str} ({record_type})\n[목표 분량]: {target_bytes} Byte\n[학생 관찰 메모]: {student_memo}"
                            if source_summary_text:
                                prompt_base += "\n\n" + source_summary_text
                            
                            response = client.chat.completions.create(
                                model=selected_model,
                                messages=[
                                    {"role": "system", "content": system_prompt},
                                    {"role": "user", "content": prompt_base}
                                ]
                            )
                            draft_res = response.choices[0].message.content.strip()
                            
                            input_tokens, output_tokens, tokens_count = extract_usage_tokens(response)
                            if tokens_count > 0:
                                st.session_state.total_tokens_used += tokens_count

                            draft_res, adj_in, adj_out, adj_total = enforce_target_bytes_with_llm(
                                client, selected_model, draft_res, target_bytes, record_type, max_attempts=1
                            )
                            input_tokens += adj_in
                            output_tokens += adj_out
                            tokens_count += adj_total
                            if adj_total > 0:
                                st.session_state.total_tokens_used += adj_total
                            
                            new_record = {
                                "info": info_str, "id_val": resolved_student_id, "name": resolved_student_name,
                                "memo": student_memo, "draft": draft_res, "record_type": record_type, "target_bytes": target_bytes,
                                "source_signature": source_signature, "draft_revision": 0
                            }
                            if tokens_count > 0:
                                new_record["tokens"] = tokens_count
                                new_record["input_tokens"] = input_tokens
                                new_record["output_tokens"] = output_tokens
                            
                            if curr_std_idx is not None and not st.session_state.is_new_student and curr_std_idx < len(st.session_state.student_records):
                                st.session_state.student_records[curr_std_idx] = new_record
                            else:
                                st.session_state.student_records.append(new_record)
                                st.session_state.current_student_idx = len(st.session_state.student_records) - 1
                                st.session_state.is_new_student = False
                                
                            st.success(f"{resolved_student_id} {resolved_student_name} 초안 생성이 완료되었습니다!")
                            st.session_state.attachments["student"] = []
                            std_key = f"uploader_std_{st.session_state.uploader_key_std}"
                            if std_key in st.session_state:
                                del st.session_state[std_key]
                            st.session_state.uploader_key_std += 1
                        except Exception as e:
                            st.session_state.is_generating_draft = False
                            st.error(f"초안 생성 중 오류 발생: {str(e)}")
                        else:
                            st.session_state.is_generating_draft = False
                            st.rerun()

        default_draft = ""
        if (not st.session_state.is_new_student) and curr_std_idx is not None and curr_std_idx < len(st.session_state.student_records):
            default_draft = st.session_state.student_records[curr_std_idx]["draft"]

        if default_draft:
            with st.container(key="draft_result"):
                st.divider()
                st.markdown("**생성된 생기부 초안 및 영역별 NEIS 검수**")
                
                # v13: 수정 후 Streamlit textarea가 이전 값을 붙잡지 않도록 초안 revision별 새 key를 사용합니다.
                curr_record_for_widget = st.session_state.student_records[curr_std_idx]
                revision = int(curr_record_for_widget.get("draft_revision", 0))
                pending_draft_key = f"_pending_draft_text_{curr_std_idx}"
                if pending_draft_key in st.session_state:
                    pending_text = st.session_state.pop(pending_draft_key)
                    revision += 1
                    curr_record_for_widget["draft_revision"] = revision
                    curr_record_for_widget["draft"] = pending_text
                    default_draft = pending_text
                ta_key = f"draft_ta_{curr_std_idx}_{revision}"
                if ta_key not in st.session_state:
                    st.session_state[ta_key] = curr_record_for_widget.get("draft", default_draft)

                edited_draft = st.text_area(
                    "초안 편집",
                    height=112,
                    label_visibility="collapsed",
                    key=ta_key
                )
                if curr_std_idx is not None and curr_std_idx < len(st.session_state.student_records):
                    st.session_state.student_records[curr_std_idx]["draft"] = edited_draft
                    
                c_cnt, b_cnt = calculate_neis_bytes(edited_draft)
                max_bytes = int(st.session_state.student_records[curr_std_idx].get("target_bytes", target_bytes))
                
                col_b1, col_b2, col_b3 = st.columns([1, 1, 2], vertical_alignment="center")
                with col_b1:
                    st.metric("글자수 (공백 포함)", f"{c_cnt} 자")
                with col_b2:
                    st.metric("NEIS 바이트", f"{b_cnt} / {max_bytes} Byte")
                with col_b3:
                    progress_ratio = b_cnt / float(max_bytes)
                    progress_val = min(1.0, progress_ratio)
                    st.write(f"[{record_type.split(' ')[0]}] 용량 달성도 (최대 {max_bytes} Byte)")
                    if b_cnt > max_bytes:
                        over_pct = int(progress_ratio * 100)
                        overflow_html = (
                            '<div style="background:#3f1d24;border-radius:8px;height:12px;overflow:hidden;border:1px solid #ef4444;">'
                            '<div style="width:100%;height:100%;background:#ef4444;"></div></div>'
                            f'<div style="color:#f87171;font-size:0.82rem;margin-top:4px;font-weight:700;">{over_pct}% · 목표 초과</div>'
                        )
                        st.markdown(overflow_html, unsafe_allow_html=True)
                    else:
                        st.progress(progress_val)
                    
                if b_cnt > max_bytes:
                    st.error(f"목표 용량({max_bytes} Byte)을 {b_cnt - max_bytes} Byte 초과하였습니다. 아래 추가 수정 요청으로 줄이거나 다시 생성해 주세요.")
                elif b_cnt < int(max_bytes * 0.8):
                    st.warning(f"현재 {b_cnt} Byte로 목표({max_bytes} Byte)의 80%보다 적습니다. 아래 추가 수정 요청에서 내용 보강을 요청할 수 있습니다.")
                else:
                    st.success(f"목표 분량 범위에 근접했습니다. 현재 {b_cnt} / {max_bytes} Byte")
                    
                forbidden_found = check_forbidden_words(edited_draft)
                if forbidden_found:
                    st.error(f"기재 주의 표현 감지: {', '.join(forbidden_found)}")
                else:
                    st.success("기재 금지어 검수 통과: 주요 기재 금지어가 감지되지 않았습니다.")

                validation_details = build_validation_issue_details(edited_draft, b_cnt, max_bytes, forbidden_found)
                if validation_details:
                    with st.expander("검수에서 확인된 문제와 수정 방법", expanded=True):
                        for issue_title, issue_reason, issue_fix in validation_details:
                            st.markdown(f"**{issue_title}**")
                            st.write(f"- 문제: {issue_reason}")
                            st.write(f"- 수정 방향: {issue_fix}")
                else:
                    st.caption("분량과 주요 기재 주의 표현 검수에서 별도 문제가 확인되지 않았습니다.")
                    
                if curr_std_idx is not None and curr_std_idx < len(st.session_state.student_records):
                    curr_record = st.session_state.student_records[curr_std_idx]
                    if "tokens" in curr_record:
                        std_in = curr_record.get("input_tokens", 0)
                        std_out = curr_record.get("output_tokens", 0)
                        if std_in or std_out:
                            st.caption(f"이번 초안 생성: 입력 **{std_in:,}** · 출력 **{std_out:,}** · 합계 **{curr_record['tokens']:,} Tokens**")
                        else:
                            st.caption(f"이번 초안 생성에 소모된 토큰: **{curr_record['tokens']:,} Tokens**")
                    if curr_record.get("refine_tokens", 0):
                        st.caption(
                            f"최근 추가 수정: 입력 **{curr_record.get('refine_input_tokens', 0):,}** · "
                            f"출력 **{curr_record.get('refine_output_tokens', 0):,}** · "
                            f"합계 **{curr_record.get('refine_tokens', 0):,} Tokens** · "
                            f"현재 **{curr_record.get('last_refine_bytes', 0):,} Byte**"
                        )

                    previous_version = curr_record.get("previous_draft", "")
                    if previous_version and previous_version != edited_draft:
                        prev_chars, prev_bytes = calculate_neis_bytes(previous_version)
                        curr_chars_cmp, curr_bytes_cmp = calculate_neis_bytes(edited_draft)
                        with st.expander("수정 전 초안과 비교", expanded=True):
                            compare_left, compare_right = st.columns(2)
                            compare_style = (
                                "background:#1e293b;border:1px solid #334155;border-radius:10px;padding:12px 14px;"
                                "margin:6px 0 4px 0;white-space:pre-wrap;word-break:break-word;line-height:1.58;"
                                "font-size:0.92rem;color:#e2e8f0;user-select:text;-webkit-user-select:text;cursor:text;"
                                "max-height:230px;overflow-y:auto;"
                            )
                            with compare_left:
                                st.markdown(f"**수정 전 ({prev_bytes:,} Byte)**")
                                prev_html = f'<div style="{compare_style}">{html.escape(previous_version.strip())}</div>'
                                st.markdown(prev_html, unsafe_allow_html=True)
                            with compare_right:
                                st.markdown(f"**현재 초안 ({curr_bytes_cmp:,} Byte)**")
                                curr_html = f'<div style="{compare_style}">{html.escape(edited_draft.strip())}</div>'
                                st.markdown(curr_html, unsafe_allow_html=True)

                    refine_key = f"draft_refine_prompt_{curr_std_idx}"
                    refine_prompt = st.text_input(
                        "추가 수정 요청",
                        key=refine_key,
                        placeholder="예: 현재 내용은 유지하되 1450~1500바이트에 맞추고, 수학적 분석 과정을 한 문장 더 구체화해줘."
                    )
                    is_refining = st.session_state.get("is_refining_draft", False)
                    if st.button("수정 중..." if is_refining else "초안 수정", type="primary", use_container_width=True, disabled=is_refining, key=f"btn_refine_{curr_std_idx}"):
                        if refine_prompt.strip():
                            st.session_state.is_refining_draft = True
                            st.session_state["_pending_refine_prompt"] = refine_prompt.strip()
                            st.session_state._run_refine_once = True
                            st.rerun()
                        else:
                            st.warning("수정할 내용을 입력해 주세요.")

                    if (st.session_state.get("is_refining_draft", False)
                            and st.session_state.get("_run_refine_once", False)
                            and st.session_state.get("_pending_refine_prompt")):
                        st.session_state._run_refine_once = False
                        client = get_openrouter_client()
                        if not client:
                            st.session_state.is_refining_draft = False
                            st.error(get_api_unavailable_message())
                        else:
                            with st.spinner(f"[{selected_model_name}] 추가 요청을 반영해 초안을 수정 중..."):
                                try:
                                    request_text = st.session_state.get("_pending_refine_prompt", "")
                                    current_bytes = calculate_neis_bytes(edited_draft)[1]
                                    byte_only_request = is_byte_adjustment_request(request_text)
                                    byte_instruction = build_byte_adjustment_instruction(current_bytes, max_bytes)

                                    refine_system = (
                                        "학교생활기록부 초안을 최소 수정합니다. 원본 첨부자료나 과거 분석은 보지 않습니다. "
                                        "현재 초안의 사실관계를 유지하고 사용자가 요청한 부분만 고치세요. 새 사실을 만들지 마세요. "
                                        "결과 문단만 출력하세요. 문장 끝은 '~함', '~임' 중심으로 유지하세요."
                                    )
                                    if byte_only_request:
                                        refine_user = (
                                            f"{byte_instruction}\n사용자 요청: {request_text}\n\n[현재 초안]\n{edited_draft}"
                                        )
                                    else:
                                        refine_user = (
                                            f"사용자 요청: {request_text}\n{byte_instruction}\n\n[현재 초안]\n{edited_draft}"
                                        )

                                    refine_response = client.chat.completions.create(
                                        model=selected_model,
                                        messages=[
                                            {"role": "system", "content": refine_system},
                                            {"role": "user", "content": refine_user}
                                        ]
                                    )
                                    revised_text = refine_response.choices[0].message.content.strip()
                                    rin, rout, rtotal = extract_usage_tokens(refine_response)
                                    if rtotal > 0:
                                        st.session_state.total_tokens_used += rtotal

                                    # v13: 분량 조절 요청은 1회 호출로 끝냄. 내용 수정일 때만 결과가 크게 벗어난 경우 1회 보정.
                                    ain = aout = atotal = 0
                                    revised_bytes = calculate_neis_bytes(revised_text)[1]
                                    if (not byte_only_request) and not (int(max_bytes * 0.88) <= revised_bytes <= max_bytes):
                                        revised_text, ain, aout, atotal = enforce_target_bytes_with_llm(
                                            client, selected_model, revised_text, max_bytes, record_type, max_attempts=1
                                        )
                                        rin += ain
                                        rout += aout
                                        rtotal += atotal
                                        if atotal > 0:
                                            st.session_state.total_tokens_used += atotal
                                    curr_record["previous_draft"] = edited_draft
                                    curr_record["draft"] = revised_text
                                    # 최초 생성 토큰과 추가 수정 토큰을 분리해 표시합니다.
                                    curr_record["refine_tokens"] = rtotal
                                    curr_record["refine_input_tokens"] = rin
                                    curr_record["refine_output_tokens"] = rout
                                    curr_record["last_refine_request"] = request_text
                                    curr_record["last_refine_bytes"] = calculate_neis_bytes(revised_text)[1]
                                    st.session_state[f"_pending_draft_text_{curr_std_idx}"] = revised_text
                                    st.session_state["_pending_refine_prompt"] = ""
                                    st.session_state.is_refining_draft = False
                                except Exception as e:
                                    st.session_state.is_refining_draft = False
                                    st.error(f"초안 수정 중 오류 발생: {str(e)}")
                                else:
                                    st.rerun()

elif mode == "생기부 검수/진단":
    # 2. Streamlit 네이티브 스크롤 컨테이너 (CSS dvh 동적 가용 높이 제어)
    with st.container(border=False, key="content_area"):
        if st.session_state.is_new_eval:
            st.caption("**현재 상태: 새 검토 중**")
        else:
            curr_eval_idx = st.session_state.current_eval_idx
            if curr_eval_idx is not None and curr_eval_idx < len(st.session_state.eval_records):
                curr_rec = st.session_state.eval_records[curr_eval_idx]
                st.caption(f"**현재 검토: {curr_rec.get('title', '이전 검토')}**")
            else:
                st.caption("**현재 상태: 새 검토 중**")

        st.markdown("**생기부 전문 진단 및 검수**")
        st.caption("기존에 작성된 생기부 문장을 업로드하거나 직접 입력하시면 지침 위반, 장단점, 문체 오류를 정밀 분석합니다.")
        
        curr_eval_idx = st.session_state.current_eval_idx
        if (not st.session_state.is_new_eval) and curr_eval_idx is not None and curr_eval_idx < len(st.session_state.eval_records):
            curr_eval_rec = st.session_state.eval_records[curr_eval_idx]
            default_eval_text = curr_eval_rec.get("target_text", "")
            default_eval_result = curr_eval_rec.get("result", "")
        else:
            default_eval_text = st.session_state.get("eval_target_text", "")
            default_eval_result = st.session_state.get("eval_result", "")
        
        col_e1, col_e2 = st.columns([3, 1])
        with col_e2:
            if default_eval_result or st.session_state.get("eval_text_widget"):
                if st.button("새로 검수하기", use_container_width=True):
                    reset_eval_form()
                    st.rerun()

        if "eval_text_widget" not in st.session_state:
            st.session_state["eval_text_widget"] = st.session_state.get("draft_eval_text", "")

        eval_input_text = st.text_area(
            "검수할 생기부 문장 직접 입력",
            height=140,
            placeholder="검수하고자 하는 생기부 특기사항 문단을 복사해서 붙여넣으세요.",
            key="eval_text_widget"
        )
        st.session_state["draft_eval_text"] = eval_input_text
        
        # 공통 첨부파일 패널 렌더링 (eval 스코프 지정 - 입력창 바로 아래 배치)
        render_attachments_panel(uploader_key=f"uploader_eval_{st.session_state.uploader_key_eval}", scope="eval")
        
        target_eval_text = eval_input_text
        
        if target_eval_text:
            c_cnt, b_cnt = calculate_neis_bytes(target_eval_text)
            st.caption(f"검수 대상 분량: {c_cnt}자 / **{b_cnt} Byte** (NEIS 기준)")
            
        is_eval_busy = st.session_state.get("is_running_eval", False)
        if st.button("검수 진행 중..." if is_eval_busy else "생기부 정밀 진단 시작", type="primary", use_container_width=True, disabled=is_eval_busy):
            if target_eval_text or any(item["type"] == "image" for item in st.session_state.attachments["eval"]) or any(item["type"] == "document" for item in st.session_state.attachments["eval"]):
                st.session_state.is_running_eval = True
                st.session_state._run_eval_once = True
                st.rerun()
            else:
                st.warning("검수할 파일이나 텍스트를 입력해 주세요.")

        if st.session_state.get("is_running_eval", False) and st.session_state.get("_run_eval_once", False):
            st.session_state._run_eval_once = False
            if not check_openrouter_model_availability(selected_model):
                st.error("현재 선택한 AI 모델을 OpenRouter에서 사용할 수 없습니다. 관리자에게 모델 설정을 확인해 주세요.")
                st.session_state.is_running_eval = False
            else:
                client = get_openrouter_client()
                if not client:
                    st.error(get_api_unavailable_message())
                    st.session_state.is_running_eval = False
                else:
                    with st.spinner(f"[{selected_model_name}] 생기부 정밀 분석 및 오류 검수 진행 중..."):
                        try:
                            eval_system_prompt = (
                                "당신은 대한민국 학교생활기록부 정밀 검수 평가관입니다.\n"
                                "제출된 생기부 텍스트 및 첨부자료를 분석하여 아래 구조에 맞춰 상세히 평가 리포트를 작성하세요:\n\n"
                                "1. 지침 위반 및 기재 금지어 적발: (대회, 수상, 외부기관, 어학성적, 사교육 유발 요소 여부 적발)\n"
                                "2. 문체 및 오탈자/비문 진단: (학생 이름/주어 시작 문장 오남용 여부, '~함/임' 어조 미준수 여부, 맞춤법 적발)\n"
                                "3. 작성의 장점 (강점): (구체적 수행과정, 도구 활용, 성취수준 표현 우수성 진술)\n"
                                "4. 작성의 단점 및 보완점: (추상적이거나 단순 총평에 그친 부분 지적)\n"
                                "5. 최종 개선/수정 제안 문장: (지침을 완벽히 준수한 최종 완성 문단 제시)\n"
                            )
                            
                            prompt_base = f"[검수 대상 생기부 텍스트]:\n{target_eval_text}"
                            user_payload = compile_api_payload(prompt_base, selected_model_name, scope="eval")
                            
                            response = client.chat.completions.create(
                                model=selected_model,
                                messages=[
                                    {"role": "system", "content": eval_system_prompt},
                                    {"role": "user", "content": user_payload}
                                ]
                            )
                            eval_res_str = response.choices[0].message.content
                            
                            input_tokens, output_tokens, tokens_count = extract_usage_tokens(response)
                            if tokens_count > 0:
                                st.session_state.total_tokens_used += tokens_count
                            
                            # 자동 제목 생성 (12~18자)
                            if target_eval_text and target_eval_text.strip():
                                clean_text = target_eval_text.strip().replace("\n", " ")
                                title_str = clean_text[:15] + "..." if len(clean_text) > 15 else clean_text
                            else:
                                title_str = "첨부자료 검토"

                            from datetime import datetime
                            now_str = datetime.now().strftime("%Y-%m-%d %H:%M")

                            if curr_eval_idx is not None and not st.session_state.is_new_eval and curr_eval_idx < len(st.session_state.eval_records):
                                old_rec = st.session_state.eval_records[curr_eval_idx]
                                created_at_val = old_rec.get("created_at", old_rec.get("updated_at", now_str))
                                new_eval_record = {
                                    "title": title_str,
                                    "target_text": target_eval_text,
                                    "result": eval_res_str,
                                    "model": selected_model_name,
                                    "tokens": tokens_count,
                                    "input_tokens": input_tokens,
                                    "output_tokens": output_tokens,
                                    "created_at": created_at_val,
                                    "updated_at": now_str
                                }
                                st.session_state.eval_records[curr_eval_idx] = new_eval_record
                            else:
                                new_eval_record = {
                                    "title": title_str,
                                    "target_text": target_eval_text,
                                    "result": eval_res_str,
                                    "model": selected_model_name,
                                    "tokens": tokens_count,
                                    "input_tokens": input_tokens,
                                    "output_tokens": output_tokens,
                                    "created_at": now_str,
                                    "updated_at": now_str
                                }
                                st.session_state.eval_records.append(new_eval_record)
                                st.session_state.current_eval_idx = len(st.session_state.eval_records) - 1
                                st.session_state.is_new_eval = False

                            st.session_state.eval_result = eval_res_str
                            st.session_state.eval_target_text = target_eval_text

                            # 성공 시 첨부 파일 리셋 및 uploader key 갱신으로 상태 완전 초기화
                            st.session_state.attachments["eval"] = []
                            eval_key = f"uploader_eval_{st.session_state.uploader_key_eval}"
                            if eval_key in st.session_state:
                                del st.session_state[eval_key]
                            st.session_state.uploader_key_eval += 1
                        except Exception as e:
                            st.session_state.is_running_eval = False
                            st.error(f"진단 중 오류 발생: {str(e)}")
                        else:
                            st.session_state.is_running_eval = False
                            st.rerun()

        if default_eval_result:
            st.divider()
            st.markdown("**AI 생기부 정밀 진단 결과**")
            
            if curr_eval_idx is not None and curr_eval_idx < len(st.session_state.eval_records):
                rec_info = st.session_state.eval_records[curr_eval_idx]
                eval_in = rec_info.get("input_tokens", 0)
                eval_out = rec_info.get("output_tokens", 0)
                if eval_in or eval_out:
                    st.caption(f"검토 모델: **{rec_info.get('model', selected_model_name)}** | 입력 **{eval_in:,}** · 출력 **{eval_out:,}** · 합계 **{rec_info.get('tokens', 0):,} Tokens**")
                else:
                    st.caption(f"검토 모델: **{rec_info.get('model', selected_model_name)}** | 사용 토큰: **{rec_info.get('tokens', 0):,} Tokens**")
                rec_date_val = rec_info.get("updated_at", rec_info.get("created_at", ""))
                rec_model_val = rec_info.get("model", selected_model_name)
            else:
                from datetime import datetime
                rec_date_val = datetime.now().strftime("%Y-%m-%d %H:%M")
                rec_model_val = selected_model_name

            # 검토 결과 렌더링 (컴팩트 스타일 컨테이너)
            with st.container(key="eval_result"):
                st.markdown(default_eval_result)

            st.markdown("---")

            # TXT 다운로드 파일 생성
            full_txt_content = (
                "Chat PSDongSung 생기부 검토 결과\n"
                "====================================\n"
                f"검토 일시: {rec_date_val}\n"
                f"사용 모델: {rec_model_val}\n"
                "====================================\n\n"
                "[검토 대상 원문]\n"
                f"{default_eval_text}\n\n"
                "[AI 정밀 진단 결과]\n"
                f"{default_eval_result}\n"
            )

            # 5번 항목 (수정 제안 문장) 분리 추출
            suggestion_text = default_eval_result
            if "5. 최종 개선" in default_eval_result:
                parts = default_eval_result.split("5. 최종 개선")
                suggestion_text = "5. 최종 개선" + parts[1]
            elif "수정 제안" in default_eval_result:
                parts = default_eval_result.split("수정 제안")
                suggestion_text = "수정 제안" + parts[1]

            from datetime import datetime
            today_date = datetime.now().strftime("%Y-%m-%d")
            clean_name = (default_eval_text.strip().replace("\n", " ")[:12]) if default_eval_text else "검토"
            clean_name = "".join([c for c in clean_name if c.isalnum() or c in (" ", "_", "-")]).strip() or "검토"
            file_name_full = f"생기부검토_{today_date}_{clean_name}.txt"
            file_name_sug = f"생기부수정제안_{today_date}_{clean_name}.txt"

            col_d1, col_d2 = st.columns(2)
            with col_d1:
                st.download_button(
                    label="전체 검토 결과 TXT 다운로드",
                    data=full_txt_content.encode("utf-8"),
                    file_name=file_name_full,
                    mime="text/plain",
                    use_container_width=True
                )
            with col_d2:
                st.download_button(
                    label="수정 제안 문장만 TXT 다운로드",
                    data=suggestion_text.encode("utf-8"),
                    file_name=file_name_sug,
                    mime="text/plain",
                    use_container_width=True
                )
